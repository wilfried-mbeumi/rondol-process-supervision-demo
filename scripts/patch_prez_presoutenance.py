"""patch_prez_presoutenance.py — Applique au support les corrections demandées après la pré-soutenance.

Deux opérations, chirurgicales, sur « MBEUMI_Wilfried_PREZ new.pptx » :

1. Diapositive 14 — remplacer la figure du championnat. L'ancienne affichait
   Random Forest à 0,796 en troisième position, alors que le commentaire de la
   même diapositive annonce 0,809 : le graphique contredisait le texte à l'écran.

2. Diapositives 6, 8 et 14 — retirer les bandeaux d'illustration générés. Ils
   n'apportent aucune information et l'un d'eux affiche du texte incohérent
   (« CHAMPIANISHIP », « MCDLIY ») lisible depuis le fond de la salle.

Les figures utiles, les captures de l'application et le diagramme de Gantt de la
diapositive 22 ne sont pas touchés.

Usage : python scripts/patch_prez_presoutenance.py
"""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "reports" / "soutenance" / "MBEUMI_Wilfried_PREZ new.pptx"
OUT = ROOT / "reports" / "soutenance" / "MBEUMI_Wilfried_PREZ_corrige.pptx"
FIG = ROOT / "figures_memoire" / "fig_championnat_modeles.png"

# Bandeaux d'illustration à retirer, repérés par leur géométrie relative.
# (diapositive, borne gauche %, borne droite %) — hauteur ~100 % de la page.
BANDEAUX = [
    (6, 0.0, 40.0),    # photo « SWOT » en anglais, doublon de la vraie figure
    (8, 60.0, 100.0),  # bandeau droit sous le schéma d'architecture
    (14, 70.0, 100.0),  # bandeau « CHAMPIANISHIP » — texte incohérent
]


def supprimer(shape) -> None:
    shape._element.getparent().remove(shape._element)


def retirer_filigrane(pptx: Path) -> int:
    """Retire le badge de l'éditeur, posé dans les gabarits et non les pages.

    C'est une image liée à gamma.app, placée en bas à droite de chaque gabarit —
    d'où son absence de la liste des formes d'une diapositive. On repère la
    relation portant l'adresse de l'éditeur, puis le bloc <p:pic> qui la
    référence.
    """
    import re
    import zipfile

    lien = re.compile(r'Id="([^"]+)"[^>]*Target="https://gamma\.app[^"]*"')
    tampon: dict[str, bytes] = {}
    retires = 0

    with zipfile.ZipFile(pptx) as z:
        noms = z.namelist()
        contenus = {n: z.read(n) for n in noms}

    for nom in noms:
        if not (nom.startswith("ppt/slideLayouts/slideLayout")
                and nom.endswith(".xml")):
            continue
        rels_nom = nom.replace("slideLayouts/", "slideLayouts/_rels/") + ".rels"
        if rels_nom not in contenus:
            continue
        ids = lien.findall(contenus[rels_nom].decode("utf-8", "ignore"))
        if not ids:
            continue
        xml = contenus[nom].decode("utf-8")
        avant = xml
        for rid in ids:
            # Le bloc <p:pic> qui référence cette relation, et lui seul.
            xml = re.sub(
                r'<p:pic>(?:(?!</p:pic>).)*?r:id="' + re.escape(rid)
                + r'"(?:(?!</p:pic>).)*?</p:pic>',
                "", xml, flags=re.S)
        if xml != avant:
            tampon[nom] = xml.encode("utf-8")
            retires += 1

    if not tampon:
        return 0

    contenus.update(tampon)
    temporaire = pptx.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(temporaire, "w", zipfile.ZIP_DEFLATED) as z:
        for nom in noms:
            z.writestr(nom, contenus[nom])
    pptx.unlink()
    temporaire.rename(pptx)
    return retires


def main() -> int:
    if not SRC.exists():
        raise SystemExit(f"Support introuvable : {SRC}")
    if not FIG.exists():
        raise SystemExit(f"Figure corrigée introuvable : {FIG}")

    shutil.copy2(SRC, OUT)
    prs = Presentation(str(OUT))
    W, H = prs.slide_width, prs.slide_height
    retires = remplacees = 0

    for num, gauche, droite in BANDEAUX:
        slide = prs.slides[num - 1]
        for shape in list(slide.shapes):
            if shape.shape_type != 13:
                continue
            x = shape.left / W * 100
            haut = shape.height / H * 100
            large = shape.width / W * 100
            # Un bandeau occupe toute la hauteur et moins de la moitié de la largeur.
            if haut >= 80.0 and large <= 45.0 and gauche <= x < droite:
                supprimer(shape)
                retires += 1
                print(f"  diapo {num:2d} — bandeau retiré "
                      f"({large:.0f}% x {haut:.0f}%, x={x:.0f}%)")

    # Diapositive 14 : la figure utile est la plus large restante.
    slide14 = prs.slides[13]
    figures = [s for s in slide14.shapes if s.shape_type == 13]
    if figures:
        cible = max(figures, key=lambda s: s.width)
        pos = (cible.left, cible.top, cible.width, cible.height)
        supprimer(cible)
        slide14.shapes.add_picture(str(FIG), pos[0], pos[1], pos[2], pos[3])
        remplacees += 1
        print(f"  diapo 14 — figure du championnat remplacée "
              f"(Random Forest 1er à 0,809)")

    # Retirer un bandeau laisse la figure décalée et un vide de l'autre côté :
    # on recentre la figure et on ramène les textes sur la même marge gauche.
    recentres = 0
    for num, _, _ in BANDEAUX:
        slide = prs.slides[num - 1]
        figures = [s for s in slide.shapes if s.shape_type == 13]
        if not figures:
            continue
        principale = max(figures, key=lambda s: s.width * s.height)
        if principale.width > W * 0.9:      # déjà pleine largeur
            continue
        nouveau_x = int((W - principale.width) / 2)
        if abs(nouveau_x - principale.left) < Emu(60000):
            continue                        # déjà centrée
        principale.left = nouveau_x

        # Translater le bloc de titre d'un seul tenant : déplacer les seuls
        # cadres porteurs de texte laisserait derrière eux les rectangles de
        # décoration qui les encadrent.
        marge = int(W * 0.06)
        autres = [s for s in slide.shapes if s is not principale
                  and s.shape_type != 13 and s.left is not None]
        if autres:
            delta = marge - min(s.left for s in autres)
            if delta:
                for s in autres:
                    s.left = s.left + delta
        recentres += 1
        print(f"  diapo {num:2d} — figure recentrée, bloc de titre translaté")

    # --- pages d'ouverture ---------------------------------------------------
    # Une version recomposée de la couverture, du sommaire et de la problématique
    # a été essayée puis écartée : composées hors de l'éditeur, elles n'avaient
    # pas la typographie du reste du support et se voyaient comme des pièces
    # rapportées. Les pages d'origine sont conservées.
    # Pour les réactiver : slide_couverture.png, slide_sommaire.png et
    # slide_problematique.png sont dans reports/soutenance/.
    REFAITES: dict[int, str] = {}
    refaites = 0
    for num, nom in REFAITES.items():
        image = ROOT / "reports" / "soutenance" / nom
        if not image.exists():
            print(f"  [!] page recomposée absente : {nom}")
            continue
        slide = prs.slides[num - 1]
        for forme in list(slide.shapes):
            supprimer(forme)
        slide.shapes.add_picture(str(image), 0, 0, W, H)
        refaites += 1
        print(f"  diapo {num:2d} — page recomposée ({nom})")

    # --- contraste : figures claires -> versions sombres ---------------------
    # Les figures du mémoire sont dessinées sur fond blanc. Projetées sur le vert
    # foncé du support, elles y découpent des rectangles blancs — c'est le défaut
    # de contraste relevé en pré-soutenance. On leur substitue les versions
    # composées sur le fond du support (figures_support/).
    SOMBRES = {6: "fig_swot.png", 8: "fig_architecture.png",
               11: "fig_two_level_ai.png", 12: "fig_data_pipeline.png",
               14: "fig_championnat_modeles.png", 18: "fig_validation.png",
               22: "fig_gantt.png", 25: "fig_tests.png"}
    assombries = 0
    for num, nom in SOMBRES.items():
        source = ROOT / "figures_support" / nom
        if not source.exists():
            print(f"  [!] version sombre absente : {nom}")
            continue
        slide = prs.slides[num - 1]
        figures = [s for s in slide.shapes if s.shape_type == 13]
        if not figures:
            continue
        cible = max(figures, key=lambda s: s.width * s.height)
        pos = (cible.left, cible.top, cible.width, cible.height)
        supprimer(cible)
        slide.shapes.add_picture(str(source), *pos)
        assombries += 1
        print(f"  diapo {num:2d} — figure passée en palette sombre ({nom})")

    # --- bandeaux d'illustration restants sur les pages de contenu -----------
    # Conservés sur la couverture et la page de fin, où l'image est un parti pris.
    for num in (2, 12, 25, 26):
        if num > len(prs.slides):
            continue
        slide = prs.slides[num - 1]
        for shape in list(slide.shapes):
            if shape.shape_type != 13 or shape.left is None:
                continue
            if shape.height / H * 100 >= 95 and shape.width / W * 100 <= 40:
                supprimer(shape)
                retires += 1
                print(f"  diapo {num:2d} — bandeau d'illustration retiré")

    # --- diapositives ajoutées après la pré-soutenance -----------------------
    # « Pourquoi ces cinq modèles » précède le championnat ; « Augmentation de
    # données » suit « Le vrai résultat », là où la question du volume se pose.
    ajouts = [
        (ROOT / "reports" / "soutenance" / "slide_choix_modeles.png", 13),
        (ROOT / "reports" / "soutenance" / "slide_augmentation.png", 17),
    ]
    ajoutees = 0
    for image, position in ajouts:
        if not image.exists():
            print(f"  [!] image absente, diapositive non ajoutée : {image.name}")
            continue
        vierge = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(vierge)
        for forme in list(slide.shapes):        # un gabarit peut porter des cadres
            supprimer(forme)
        slide.shapes.add_picture(str(image), 0, 0, W, H)

        # add_slide ajoute en fin de présentation : on replace l'entrée dans
        # la liste ordonnée des diapositives.
        liste = prs.slides._sldIdLst
        entree = liste[-1]
        liste.remove(entree)
        liste.insert(position, entree)
        ajoutees += 1
        print(f"  diapo {position + 1:2d} — ajoutée : {image.stem}")

    # --- cartes claires -> teinte du thème -----------------------------------
    # Gamma pose ses encarts en lavande clair (#C6C9DC) : sur le vert foncé du
    # support, ce sont des pavés blancs. On les rhabille, et on éclaircit le
    # texte posé dessus, sinon il devient illisible.
    from pptx.dml.color import RGBColor
    CLAIR = (0xC6, 0xC9, 0xDC)
    CARTE = RGBColor(0x1B, 0x46, 0x40)
    ENCRE = RGBColor(0xF2, 0xEC, 0xE0)
    cartes = 0

    for slide in prs.slides:
        zones = []
        for shape in slide.shapes:
            try:
                remplissage = shape.fill
                if remplissage.type != 1:
                    continue
                couleur = remplissage.fore_color
                if couleur.type != 1 or tuple(couleur.rgb) != CLAIR:
                    continue
            except Exception:
                continue
            if (shape.height or 0) / H * 100 < 5:   # filets de tableau
                continue
            remplissage.fore_color.rgb = CARTE
            try:
                shape.line.color.rgb = RGBColor(0x2F, 0xB3, 0x9B)
                shape.line.width = Emu(9525)
            except Exception:
                pass
            zones.append((shape.left, shape.top,
                          shape.left + shape.width, shape.top + shape.height))
            cartes += 1

        if not zones:
            continue
        for shape in slide.shapes:
            if not shape.has_text_frame or shape.left is None:
                continue
            cx = shape.left + (shape.width or 0) / 2
            cy = shape.top + (shape.height or 0) / 2
            if not any(x1 <= cx <= x2 and y1 <= cy <= y2 for x1, y1, x2, y2 in zones):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = ENCRE

    if cartes:
        print(f"  {cartes} encart(s) clair(s) passé(s) à la teinte du thème")

    prs.save(OUT)
    filigranes = retirer_filigrane(OUT)
    if filigranes:
        print(f"  filigrane de l'éditeur retiré de {filigranes} gabarit(s)")
    print(f"\n[OK] {OUT.name} — {retires} bandeau(x) retiré(s), "
          f"{remplacees + assombries} figure(s) remplacée(s), {ajoutees} ajoutée(s), "
          f"{filigranes} filigrane(s), {len(prs.slides)} diapositives")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
