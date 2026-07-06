# -*- coding: utf-8 -*-
"""
Support de soutenance Rondol (15 slides, 16:9) + script oral FR.
- PPTX premium cohérent avec le mémoire (bleu pétrole / vert Rondol / blanc)
- Captures réelles + schémas + graphiques
- Notes orateur sous chaque slide
- Script oral slide par slide -> Script_Soutenance_...docx (≈ 10 min)
Sortie : reports/Soutenance_...FINAL_PRO.pptx + reports/Script_Soutenance_...docx
"""
from __future__ import annotations
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports"
ASSETS = REPORTS / "assets"
FIG = REPORTS / "memoire_figures"
CAP = REPORTS / "memoire_captures"
V2 = REPORTS / "poster_abstract" / "figures" / "generated_v2"
PPTX = REPORTS / "Soutenance_Rondol_Wilfried_Galtier_MBEUMI_FINAL_PRO.pptx"
SCRIPT_DOCX = REPORTS / "Script_Soutenance_Rondol_Wilfried_Galtier_MBEUMI.docx"

NEXA = ASSETS / "nexa_logo.png"
RONDOL = ASSETS / "rondol_logo.png"

# Palette
NAVY = RGBColor(0x1F, 0x4E, 0x79)
NAVY_D = RGBColor(0x16, 0x39, 0x5B)
BLUE2 = RGBColor(0x2E, 0x75, 0xB6)
GREEN = RGBColor(0x1B, 0x7A, 0x3D)
GREEN2 = RGBColor(0x35, 0xA8, 0x60)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x26, 0x32, 0x3A)
GREY = RGBColor(0x5A, 0x6A, 0x75)
RED = RGBColor(0xB0, 0x3A, 0x2E)
AMBER = RGBColor(0xB7, 0x79, 0x1F)

SW, SH = 13.333, 7.5
HEAD = "Calibri"
BODY = "Calibri"

SCRIPT = []  # (titre, discours) pour le DOCX


# --------------------------------------------------------------------------- #
def prs_new():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill, line=None, lw=1.0, rounded=False, shadow=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line; shape.line.width = Pt(lw)
    shape.shadow.inherit = False
    return shape


def text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if p.get("space_before") is not None:
            para.space_before = Pt(p["space_before"])
        para.space_after = Pt(p.get("space_after", 4))
        if p.get("line") is not None:
            para.line_spacing = p["line"]
        segs = p["runs"] if "runs" in p else [(p.get("text", ""), p)]
        for seg in segs:
            if isinstance(seg, tuple):
                txt, st = seg
            else:
                txt, st = seg, p
            r = para.add_run(); r.text = txt
            f = r.font
            f.size = Pt(st.get("size", p.get("size", 16)))
            f.bold = st.get("bold", p.get("bold", False))
            f.italic = st.get("italic", p.get("italic", False))
            f.name = st.get("font", HEAD)
            f.color.rgb = st.get("color", p.get("color", INK))
    return tb


def image_fit(slide, path, x, y, w, h):
    iw, ih = Image.open(path).size
    br, ir = w / h, iw / ih
    if ir > br:
        nw, nh = w, w / ir
    else:
        nh, nw = h, h * ir
    nx, ny = x + (w - nw) / 2, y + (h - nh) / 2
    return slide.shapes.add_picture(str(path), Inches(nx), Inches(ny), Inches(nw), Inches(nh))


def logo_chip(slide, path, x, y, h, chip=True):
    iw, ih = Image.open(path).size
    w = h * iw / ih
    if chip:
        rect(slide, x - 0.12, y - 0.1, w + 0.24, h + 0.2, WHITE, rounded=True)
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    return w


def top_band(slide, title, idx):
    rect(slide, 0, 0, SW, 1.02, NAVY)
    rect(slide, 0, 0, 0.22, 1.02, GREEN2)            # motif : chip vert à gauche
    text(slide, 0.55, 0.12, 10.5, 0.8,
         [{"text": title, "size": 27, "bold": True, "color": WHITE}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(slide, SW - 1.4, 0.12, 0.9, 0.8,
         [{"text": f"{idx:02d}", "size": 18, "bold": True, "color": RGBColor(0x9D, 0xC3, 0xE6),
           "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)


def footer(slide):
    text(slide, 0.55, SH - 0.42, 9.5, 0.34,
         [{"text": "Wilfried Galtier MBEUMI  ·  Nexa Digital School × Rondol Industrie  ·  2025–2026",
           "size": 9, "color": GREY}], anchor=MSO_ANCHOR.MIDDLE)
    try:
        logo_chip(slide, RONDOL, SW - 1.55, SH - 0.46, 0.30, chip=False)
    except Exception:
        pass


def card(slide, x, y, w, h, title, lines, accent=BLUE2, title_size=15, body_size=12.5):
    rect(slide, x, y, w, h, WHITE, line=RGBColor(0xDD, 0xE3, 0xE8), lw=1, rounded=True)
    rect(slide, x, y, 0.12, h, accent, rounded=False)
    paras = [{"text": title, "size": title_size, "bold": True, "color": accent, "space_after": 5}]
    for ln in lines:
        paras.append({"runs": [("▪  ", {"size": body_size, "color": accent, "bold": True}),
                               (ln, {"size": body_size, "color": INK})], "space_after": 3, "line": 1.0})
    text(slide, x + 0.28, y + 0.16, w - 0.42, h - 0.3, paras)


def stat(slide, x, y, w, h, value, label, accent=NAVY):
    rect(slide, x, y, w, h, WHITE, line=RGBColor(0xDD, 0xE3, 0xE8), lw=1, rounded=True)
    text(slide, x + 0.1, y + 0.18, w - 0.2, h * 0.55,
         [{"text": value, "size": 30, "bold": True, "color": accent, "align": PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(slide, x + 0.1, y + h * 0.6, w - 0.2, h * 0.36,
         [{"text": label, "size": 11.5, "color": GREY, "align": PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.TOP)


def notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt


def caption(slide, x, y, w, txt):
    text(slide, x, y, w, 0.34,
         [{"text": txt, "size": 10.5, "italic": True, "color": GREY, "align": PP_ALIGN.CENTER}])



def _overview(prs):
    """Slide bilan : tout le projet compris en un coup d'oeil."""
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT)
    rect(s, 0, 0, SW, 1.02, NAVY); rect(s, 0, 0, 0.22, 1.02, GREEN2)
    text(s, 0.55, 0.1, SW - 1.1, 0.82,
         [{"text": "PROJET EN UN COUP D'OEIL", "size": 22, "bold": True, "color": WHITE}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.6, 1.2, SW - 1.2, 0.5,
         [{"text": "De la contrainte industrielle a l'aide a la decision deployee",
           "size": 14, "italic": True, "color": GREY}])
    blocks = [
        ("1 - PROBLEME", ["Extrusion empirique", "Essais couteux,", "non reproductibles"], RED),
        ("2 - DONNEES", ["8 essais reels (12 capteurs)", "+ augmentation simulee", "documentee depuis l'echantillon"], BLUE2),
        ("3 - MODELE IA", ["Championnat 5 modeles", "RandomForest retenu", "F1-macro 0,918 (essai reel)"], GREEN),
        ("4 - APPLICATION", ["Streamlit : score stabilite", "alertes + recommandations", "(regles expertes tracables)"], NAVY),
        ("5 - VALEUR", ["Anticipe l'instabilite", "oriente les reglages", "decision finale humaine"], GREEN2),
    ]
    bw, bh, y0 = 2.28, 2.15, 2.05
    xs = [0.5 + i * (bw + 0.18) for i in range(5)]
    for (title, lines, col), x in zip(blocks, xs):
        rect(s, x, y0, bw, bh, WHITE, line=RGBColor(0xDD, 0xE3, 0xE8), lw=1, rounded=True, shadow=True)
        rect(s, x, y0, bw, 0.42, col, rounded=False)
        text(s, x + 0.12, y0 + 0.03, bw - 0.24, 0.38,
             [{"text": title, "size": 11.5, "bold": True, "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + 0.14, y0 + 0.55, bw - 0.28, bh - 0.65,
             [{"text": ln, "size": 11, "color": INK, "space_after": 2, "line": 1.02} for ln in lines])
    for i in range(4):
        text(s, xs[i] + bw - 0.04, y0 + bh / 2 - 0.25, 0.26, 0.5,
             [{"text": ">", "size": 18, "bold": True, "color": GREEN2, "align": PP_ALIGN.CENTER}],
             anchor=MSO_ANCHOR.MIDDLE)
    stats = [("12", "Capteurs de temperature", NAVY), ("800", "Fenetres simulees (aug.)", BLUE2),
             ("0,918", "F1-macro RF (essai reel)", GREEN), ("693", "Tests automatises OK", GREEN2)]
    sw2, sh2, sy = 2.9, 1.5, 4.75
    for (v, l, c), x in zip(stats, [0.5 + i * (sw2 + 0.18) for i in range(4)]):
        stat(s, x, sy, sw2, sh2, v, l, accent=c)
    caption(s, 0.5, 6.5, SW - 1.0,
            "Logique metier d'extrusion reelle + modele supervise honnetement valide + agent explicable - pas une IA decorative.")
    notes(s, "En une slide, tout mon projet. Le probleme : l'extrusion de composants de batteries reste "
              "empirique et couteuse. Les donnees : huit essais reels, augmentes par un dataset simule "
              "documente. Le modele : un championnat de cinq algorithmes dont le RandomForest ressort a "
              "0,918 en validation par essai reel. L'application : un outil Streamlit qui affiche un score "
              "de stabilite et des recommandations issues de regles expertes. La valeur : anticiper "
              "l'instabilite et orienter les reglages, la decision restant humaine.")


# --------------------------------------------------------------------------- #
def build(prs):
    # ---- S1 Titre (fond bleu) ------------------------------------------- #
    s = blank(prs)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 0, 0.28, SH, GREEN2)
    logo_chip(s, NEXA, 0.7, 0.55, 0.7)
    iw, ih = Image.open(RONDOL).size
    rw = 0.62 * iw / ih
    logo_chip(s, RONDOL, SW - 0.7 - rw, 0.58, 0.62)
    text(s, 1.2, 2.0, SW - 2.4, 0.5,
         [{"text": "MÉMOIRE DE THÈSE PROFESSIONNELLE — SOUTENANCE", "size": 15, "bold": True,
           "color": GREEN2, "align": PP_ALIGN.CENTER}])
    text(s, 1.1, 2.55, SW - 2.2, 2.0, [
        {"text": "Système d'IA prédictif d'aide à la décision pour l'extrusion bivis",
         "size": 30, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "line": 1.05},
        {"text": "de composants de batteries tout-solide (dry / semi-dry)",
         "size": 30, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "line": 1.05},
    ])
    text(s, 1.5, 4.5, SW - 3.0, 0.5,
         [{"text": "Un prototype professionnel de jumeau numérique — extrudeuse bivis 10,5 mm Rondol Industrie",
           "size": 14, "italic": True, "color": RGBColor(0xCC, 0xE0, 0xD4), "align": PP_ALIGN.CENTER}])
    text(s, 1.5, 5.35, SW - 3.0, 0.45,
         [{"text": "Présenté et soutenu par  Wilfried Galtier MBEUMI", "size": 17, "bold": True,
           "color": WHITE, "align": PP_ALIGN.CENTER}])
    text(s, 1.5, 5.95, SW - 3.0, 0.8, [
        {"text": "Nexa Digital School — Mastère Data & Intelligence Artificielle (RNCP 37137, Niveau 7)",
         "size": 12.5, "color": RGBColor(0xCF, 0xDD, 0xEC), "align": PP_ALIGN.CENTER, "space_after": 2},
        {"runs": [("Tuteur industriel : ", {"size": 12.5, "color": RGBColor(0x9D, 0xC3, 0xE6)}),
                  ("M. Maël Gallas (Rondol Industrie)", {"size": 12.5, "bold": True, "color": WHITE}),
                  ("      Référent pédagogique : ", {"size": 12.5, "color": RGBColor(0x9D, 0xC3, 0xE6)}),
                  ("M. Moussa NDIAYE (Nexa)", {"size": 12.5, "bold": True, "color": WHITE})],
         "align": PP_ALIGN.CENTER}])
    text(s, 1.5, 6.85, SW - 3.0, 0.4,
         [{"text": "Année universitaire 2025 – 2026", "size": 12, "italic": True,
           "color": RGBColor(0x9D, 0xC3, 0xE6), "align": PP_ALIGN.CENTER}])
    notes(s, "Bonjour. Je suis Wilfried Galtier MBEUMI. Je vous présente mon mémoire de thèse "
              "professionnelle, réalisé chez Rondol Industrie : un système d'IA prédictif d'aide "
              "à la décision pour l'extrusion bivis de composants de batteries tout-solide. "
              "Je remercie M. Maël Gallas, mon tuteur industriel, et M. Moussa NDIAYE, mon référent Nexa.")
    SCRIPT.append(("Slide 1 — Page de titre (≈ 30 s)",
        "Bonjour à toutes et à tous. Je m'appelle Wilfried Galtier MBEUMI et je vous présente "
        "aujourd'hui mon mémoire de thèse professionnelle, mené au sein de Rondol Industrie dans le "
        "cadre du Mastère Data & Intelligence Artificielle de Nexa Digital School. Le sujet : la "
        "conception et le déploiement d'un système d'intelligence artificielle prédictif d'aide à la "
        "décision pour l'optimisation de l'extrusion bivis de composants de batteries tout-solide. "
        "Je remercie mon tuteur industriel M. Maël Gallas et mon référent pédagogique M. Moussa NDIAYE."))

    _overview(prs)

    # ---- S2 Contexte industriel ----------------------------------------- #
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT); top_band(s, "Contexte industriel", 2)
    card(s, 0.55, 1.35, 6.1, 1.15, "Une rupture : les batteries tout-solide",
         ["Sécurité accrue, densité d'énergie supérieure au lithium-ion"], accent=NAVY, title_size=15)
    card(s, 0.55, 2.65, 6.1, 1.15, "Un verrou : la mise en forme",
         ["Procédés humides solvantés de plus en plus contestés"], accent=BLUE2, title_size=15)
    card(s, 0.55, 3.95, 6.1, 1.15, "Une pression réglementaire",
         ["Restriction PFAS (ECHA 2023) visant le liant PVDF"], accent=AMBER, title_size=15)
    card(s, 0.55, 5.25, 6.1, 1.3, "Une réponse : l'extrusion bivis (HME)",
         ["Voie continue sèche / semi-dry — cœur de métier Rondol"], accent=GREEN, title_size=15)
    image_fit(s, FIG / "fig_domain_intersection.png", 6.95, 1.45, 6.0, 5.0)
    footer(s)
    notes(s, "Les batteries tout-solide promettent plus de sécurité et d'énergie, mais butent sur la "
              "mise en forme des matériaux. Les procédés humides solvantés et le PVDF sont sous pression "
              "réglementaire (PFAS). L'extrusion bivis à chaud, sans solvant, est le cœur de métier de "
              "Rondol — et une voie crédible. Or l'intersection extrusion + IA + batteries reste peu étudiée.")
    SCRIPT.append(("Slide 2 — Contexte industriel (≈ 45 s)",
        "Les batteries tout-solide représentent une rupture : plus sûres, plus denses en énergie. Mais "
        "leur industrialisation bute sur un verrou souvent sous-estimé — la mise en forme des électrodes "
        "et électrolytes — pour laquelle les procédés humides solvantés sont de plus en plus contestés, "
        "notamment depuis la proposition de restriction des PFAS qui vise le PVDF. L'extrusion bivis à "
        "chaud, sèche ou semi-sèche, apporte une réponse : c'est précisément le savoir-faire de Rondol "
        "Industrie. Comme le montre ce schéma, l'intersection entre extrusion bivis, intelligence "
        "artificielle et batteries tout-solide reste très peu couverte par la littérature : c'est là que "
        "se situe mon projet."))

    # ---- S3 Problématique (fond bleu) ----------------------------------- #
    s = blank(prs); rect(s, 0, 0, SW, SH, NAVY); rect(s, 0, 0, 0.28, SH, GREEN2)
    text(s, 1.0, 1.1, SW - 2.0, 0.6,
         [{"text": "PROBLÉMATIQUE", "size": 20, "bold": True, "color": GREEN2, "align": PP_ALIGN.CENTER}])
    rect(s, 1.5, 2.2, SW - 3.0, 3.1, NAVY_D, line=RGBColor(0x9D, 0xC3, 0xE6), lw=1.5, rounded=True)
    text(s, 2.0, 2.5, SW - 4.0, 2.5,
         [{"text": "« Comment concevoir et déployer un système d'intelligence artificielle prédictif "
                   "permettant d'optimiser les paramètres d'extrusion pour la fabrication de composants "
                   "de batteries tout-solide (SSB) dry/semi-dry, afin d'améliorer la performance technique "
                   "et la compétitivité stratégique de Rondol Industrie ? »",
           "size": 21, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "line": 1.2}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.5, 5.7, SW - 3.0, 0.6,
         [{"text": "Problématique formulée et validée avec le tuteur industriel, M. Maël Gallas",
           "size": 13, "italic": True, "color": RGBColor(0xCC, 0xE0, 0xD4), "align": PP_ALIGN.CENTER}])
    notes(s, "Voici la problématique, validée avec M. Maël Gallas. Trois exigences : scientifique "
              "(prédire à partir de données réelles peu nombreuses), ingénierie (une plateforme cohérente "
              "et démontrable) et stratégique (un facteur de différenciation pour Rondol).")
    SCRIPT.append(("Slide 3 — Problématique (≈ 35 s)",
        "Ma problématique, formulée et validée avec mon tuteur industriel, est la suivante : comment "
        "concevoir et déployer un système d'IA prédictif d'aide à la décision pour optimiser les "
        "paramètres d'extrusion de composants de batteries tout-solide, afin d'améliorer la performance "
        "technique et la compétitivité de Rondol ? Elle articule trois exigences : scientifique, "
        "d'ingénierie logicielle, et stratégique."))

    # ---- S4 Objectifs --------------------------------------------------- #
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT); top_band(s, "Objectifs du projet", 4)
    card(s, 0.55, 1.45, 5.95, 4.9, "Objectifs fonctionnels", [
        "Configurer le procédé (profil de vis, thermique, dosage)",
        "Calculer les indicateurs (remplissage, résidence, volumes)",
        "Évaluer la stabilité : score + probabilité de dérive",
        "Émettre alertes hiérarchisées et recommandations chiffrées",
        "Conserver l'historique des configurations validées",
    ], accent=NAVY, title_size=18, body_size=14)
    card(s, 6.8, 1.45, 5.95, 4.9, "Objectifs techniques", [
        "Architecture en couches, calculs purs et testables",
        "Source unique de vérité : snapshot persistant",
        "Persistance durable et auto-réparatrice (Supabase)",
        "685 tests automatisés · internationalisation FR/EN",
        "Deux niveaux d'IA distingués (référence vs intégré)",
    ], accent=GREEN, title_size=18, body_size=14)
    footer(s)
    notes(s, "Côté fonctionnel : configurer, calculer, évaluer la stabilité, recommander, historiser. "
              "Côté technique : une architecture en couches testable, une source unique de vérité "
              "persistante et auto-réparatrice, 685 tests, le bilinguisme, et surtout la distinction "
              "explicite de deux niveaux d'IA.")
    SCRIPT.append(("Slide 4 — Objectifs (≈ 40 s)",
        "Les objectifs se déclinent en deux volets. Fonctionnellement, l'outil doit permettre de "
        "configurer le procédé, de calculer les indicateurs clés — taux de remplissage, temps de "
        "résidence, volumes — d'évaluer la stabilité par un score et une probabilité de dérive, "
        "d'émettre des alertes et des recommandations chiffrées, et de conserver l'historique. "
        "Techniquement, je visais une architecture en couches testable, une source unique de vérité "
        "persistante et auto-réparatrice, une suite de 685 tests, le bilinguisme, et la distinction "
        "claire entre un modèle de référence et la logique réellement intégrée."))

    # ---- S5 Données ----------------------------------------------------- #
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT); top_band(s, "Données utilisées", 5)
    stat(s, 0.55, 1.5, 2.9, 1.7, "Avril 2026", "Campagne d'essais Rondol", NAVY)
    stat(s, 3.65, 1.5, 2.9, 1.7, "12", "capteurs de température", BLUE2)
    stat(s, 0.55, 3.4, 2.9, 1.7, "8 / 11", "essais exploitables", GREEN)
    stat(s, 3.65, 3.4, 2.9, 1.7, "~50 145", "enregistrements (filière)", AMBER)
    text(s, 0.55, 5.35, 6.0, 1.4, [
        {"runs": [("Données réelles, non simulées", {"size": 13.5, "bold": True, "color": GREEN}),
                  (" — mais volume modeste : prudence statistique assumée.",
                   {"size": 13.5, "color": INK})], "line": 1.05},
        {"text": "Z1–Z8 (fourreau), filière DIE, 3 points ligne de film.",
         "size": 12.5, "color": GREY, "space_before": 4}])
    image_fit(s, FIG / "fig_data_pipeline.png", 6.85, 2.0, 6.1, 3.4)
    caption(s, 6.85, 5.45, 6.1, "Pipeline de données : du CSV brut au jeu d'apprentissage")
    footer(s)
    notes(s, "Les données viennent d'une campagne réelle d'avril 2026 : 12 capteurs de température, "
              "11 essais dont 8 exploitables, échantillonnage irrégulier — le seul capteur de filière "
              "totalise plus de 50 000 enregistrements. Données réelles mais peu nombreuses : c'est une "
              "limite que j'assume et qui oriente toute la méthodologie de validation.")
    SCRIPT.append(("Slide 5 — Données utilisées (≈ 45 s)",
        "Les données proviennent d'une campagne d'essais réelle menée chez Rondol du 7 au 13 avril 2026. "
        "Douze capteurs de température : les huit zones du fourreau, la filière, et trois points sur la "
        "ligne de film. Onze essais, dont huit exploitables après filtrage. L'échantillonnage est "
        "irrégulier, de l'ordre de la seconde à la quinzaine de secondes ; le seul capteur de filière "
        "totalise plus de cinquante mille enregistrements. Le point important : ce sont des données "
        "réelles, mais peu nombreuses. Cette rareté est assumée et conditionne toute ma validation. "
        "Le pipeline à droite résume le passage du CSV brut au jeu d'apprentissage."))

    # ---- S6 Méthodologie ------------------------------------------------ #
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT); top_band(s, "Méthodologie projet", 6)
    text(s, 0.55, 1.5, 5.6, 5.0, [
        {"text": "CRISP-DM adapté", "size": 17, "bold": True, "color": NAVY, "space_after": 8},
        {"runs": [("▪  ", {"color": GREEN, "bold": True}), ("Compréhension métier & données", {"color": INK})], "size": 14, "space_after": 6},
        {"runs": [("▪  ", {"color": GREEN, "bold": True}), ("Préparation des données (nettoyage, fenêtres)", {"color": INK})], "size": 14, "space_after": 6},
        {"runs": [("▪  ", {"color": GREEN, "bold": True}), ("Modélisation (RF / XGBoost / SVM)", {"color": INK})], "size": 14, "space_after": 6},
        {"runs": [("▪  ", {"color": GREEN, "bold": True}), ("Évaluation stricte (séparation par essai)", {"color": INK})], "size": 14, "space_after": 6},
        {"runs": [("▪  ", {"color": GREEN, "bold": True}), ("Déploiement (Streamlit Cloud + Supabase)", {"color": INK})], "size": 14, "space_after": 10},
        {"runs": [("Mode « rival » : ", {"size": 13.5, "bold": True, "color": NAVY}),
                  ("chaque hypothèse challengée plutôt que validée par défaut.",
                   {"size": 13.5, "italic": True, "color": GREY})], "line": 1.05},
    ])
    image_fit(s, FIG / "fig_crispdm.png", 6.3, 1.35, 6.6, 5.2)
    footer(s)
    notes(s, "J'ai suivi une démarche CRISP-DM adaptée à un projet mené par un développeur unique : "
              "compréhension métier et données, préparation, modélisation, évaluation, déploiement, avec "
              "des itérations. Le tuteur a imposé un mode 'rival' : challenger chaque hypothèse plutôt "
              "que de la valider par défaut — un garde-fou méthodologique permanent.")
    SCRIPT.append(("Slide 6 — Méthodologie (≈ 35 s)",
        "La conduite du projet s'appuie sur la méthodologie CRISP-DM, adaptée à un projet mené par un "
        "acteur unique : compréhension métier puis des données, préparation, modélisation, évaluation, "
        "déploiement, avec des allers-retours entre étapes. Mon tuteur a instauré un « mode rival » : "
        "exiger un contre-argument rationnel plutôt qu'une validation par défaut. Ce garde-fou a "
        "directement structuré la rigueur du travail."))

    # ---- S7 Architecture ------------------------------------------------ #
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT); top_band(s, "Architecture de la solution", 7)
    image_fit(s, FIG / "fig_architecture.png", 0.7, 1.3, 8.5, 5.3)
    text(s, 9.35, 1.6, 3.5, 5.0, [
        {"text": "Principe", "size": 16, "bold": True, "color": NAVY, "space_after": 6},
        {"text": "Dépendances strictement descendantes.", "size": 12.5, "color": INK, "space_after": 8},
        {"runs": [("« Envelopper, ", {"size": 12.5, "bold": True, "color": GREEN}),
                  ("ne pas recalculer » ", {"size": 12.5, "bold": True, "color": GREEN}),
                  ("— Network 7 appelé une seule fois.", {"size": 12.5, "color": INK})], "space_after": 8, "line": 1.05},
        {"text": "L'interface ne fait que restituer un état déjà calculé.",
         "size": 12.5, "color": GREY, "italic": True, "line": 1.05},
    ])
    footer(s)
    notes(s, "La plateforme est stratifiée : un backbone géométrique (screw_logic, Network 7), des "
              "packages purs, un moteur d'enveloppement, l'interface Streamlit, et la persistance "
              "Supabase. Principe fondateur : envelopper, ne pas recalculer. Network 7 n'est appelé "
              "qu'une fois ; l'interface ne fait que restituer un état déjà calculé. C'est ce qui "
              "garantit la testabilité et la cohérence scientifique.")
    SCRIPT.append(("Slide 7 — Architecture (≈ 45 s)",
        "La solution repose sur une architecture en couches strictes. À la base, le backbone géométrique "
        "— le module screw_logic et son calcul « Network 7 » — seule source des grandeurs procédé. "
        "Au-dessus, des packages purs de machine, matériaux et physique ; puis un moteur d'enveloppement "
        "qui enrichit l'état sans jamais le recalculer ; l'interface Streamlit ; et la persistance "
        "Supabase. Le principe fondateur est « envelopper, ne pas recalculer » : Network 7 est appelé "
        "exactement une fois, et l'interface se contente de restituer un état déjà calculé. Les "
        "dépendances sont strictement descendantes, ce qui rend le cœur scientifique testable isolément."))

    # ---- S8 Modélisation ML --------------------------------------------- #
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT); top_band(s, "Modélisation : championnat & augmentation", 8)
    image_fit(s, FIG / "fig_two_level_ai.png", 0.55, 1.3, 7.5, 4.4)
    stat(s, 8.35, 1.45, 2.2, 1.55, "0,917", "F1-macro RF (test)", NAVY)
    stat(s, 10.65, 1.45, 2.2, 1.55, "0,918", "F1-macro RF (augmenté, essai réel)", GREEN)
    card(s, 8.35, 3.2, 4.5, 2.5, "Championnat + augmentation", [
        "5 modèles, validation par essai réel non vu",
        "Sans augmentation : ~0,80 (8 essais limitent)",
        "Augmentation documentée depuis l'échantillon",
        "RandomForest retenu : 0,918 ± 0,054 → déployé",
    ], accent=GREEN, title_size=14.5, body_size=12)
    footer(s)
    notes(s, "Je compare cinq modèles supervisés en validation par essai réel non vu — le protocole "
              "honnête. Sans plus de données, tous plafonnent autour de 0,80 : ce sont les huit essais qui "
              "limitent, pas l'algorithme. Rondol n'ayant pas de base d'essais, j'ai généré un dataset simulé "
              "documenté à partir de l'échantillon — jamais aléatoirement, en reproduisant même les valeurs "
              "manquantes. Résultat : le Random Forest passe à 0,918 sur essais réels non vus et devient le "
              "modèle retenu et déployé. Le modèle prédit la stabilité ; les règles expertes recommandent.")
    SCRIPT.append(("Slide 8 — Modélisation ML (≈ 55 s)",
        "Un parti pris fort traverse le projet : distinguer deux niveaux d'intelligence artificielle. "
        "Le premier est un modèle de référence hors-ligne, le Random Forest, qui établit le potentiel "
        "prédictif : 95 % d'exactitude, 0,917 de F1-macro en test. Mais — et c'est essentiel — sous "
        "validation stricte par essai, en Leave-One-Group-Out, ce F1 réaliste tombe à 0,77, avec un "
        "écart-type de 0,11. Cet écart d'une quinzaine de points par rapport au split aléatoire vient de "
        "l'autocorrélation temporelle : je ne le dissimule pas, je le revendique comme la preuve de la "
        "rigueur de validation. Face à ce déficit, j'ai mis en œuvre une augmentation de données "
        "documentée, générée à partir de l'échantillon réel : elle porte le Random Forest à 0,918 sur "
        "essais réels non vus et divise sa variance par plus de trois. C'est le modèle retenu et déployé. "
        "Enfin, la décision n'est jamais confiée au seul modèle : il prédit la stabilité, les règles "
        "expertes traçables produisent les recommandations, et l'ingénieur tranche."))

    # ---- S9 Agent d'aide à la décision ---------------------------------- #
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT); top_band(s, "Agent d'aide à la décision", 9)
    text(s, 0.55, 1.5, 5.7, 5.0, [
        {"text": "Explicable par construction", "size": 16, "bold": True, "color": NAVY, "space_after": 8},
        {"runs": [("▪  ", {"color": GREEN, "bold": True}), ("10 règles métier (R1–R10)", {"color": INK})], "size": 14, "space_after": 6},
        {"runs": [("▪  ", {"color": GREEN, "bold": True}), ("Alertes hiérarchisées : CRITICAL / WARNING / INFO", {"color": INK})], "size": 14, "space_after": 6},
        {"runs": [("▪  ", {"color": GREEN, "bold": True}), ("Score de risque 0–100 → STABLE / SURVEILLER / CRITIQUE", {"color": INK})], "size": 14, "space_after": 6},
        {"runs": [("▪  ", {"color": GREEN, "bold": True}), ("Recommandations chiffrées avant → après", {"color": INK})], "size": 14, "space_after": 6},
        {"runs": [("▪  ", {"color": GREEN, "bold": True}), ("Traçabilité : chaque reco liée à son alerte", {"color": INK})], "size": 14, "space_after": 6},
        {"runs": [("▪  ", {"color": GREEN, "bold": True}), ("Équation thermique imposée par le tuteur", {"color": INK})], "size": 14, "space_after": 6},
    ])
    image_fit(s, V2 / "figv2_04_agent_recommendation_panel.png", 6.35, 1.55, 6.5, 4.4)
    caption(s, 6.35, 5.9, 6.5, "Panneau de recommandation de l'agent (cas C4)")
    footer(s)
    notes(s, "L'agent est explicable par construction : dix règles métier, des alertes hiérarchisées, un "
              "score de 0 à 100 traduit en trois états, et des recommandations chiffrées avant→après. "
              "Chaque recommandation est liée à l'alerte qui la motive — aucune n'est orpheline. "
              "L'équation thermique a été imposée par le tuteur industriel. C'est cette explicabilité qui "
              "conditionne l'adoption par des ingénieurs.")
    SCRIPT.append(("Slide 9 — Agent d'aide à la décision (≈ 45 s)",
        "La valeur décisionnelle vient de l'agent, explicable par construction. Il repose sur dix règles "
        "métier, produit des alertes hiérarchisées — critique, avertissement, information — et un score "
        "de risque de 0 à 100 traduit en trois états : stable, à surveiller, critique. Surtout, chaque "
        "recommandation est chiffrée, du type « ramener la teneur de 35 % à 17–20 % », et reliée "
        "explicitement à l'alerte qui la motive : aucune recommandation n'est orpheline. L'équation "
        "thermique au cœur du modèle a été imposée par mon tuteur. C'est cette traçabilité de bout en "
        "bout qui rend l'outil digne de confiance pour un ingénieur."))

    # ---- S10 Interface (grille 6 captures) ------------------------------ #
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT); top_band(s, "Interface développée — 6 pages", 10)
    grid = [
        (CAP / "cap_supervision.png", "Supervision"),
        (CAP / "cap_profile.png", "Profile"),
        (CAP / "cap_settings.png", "Settings"),
        (CAP / "cap_run_analysis.png", "Run Analysis"),
        (CAP / "cap_history.png", "History"),
        (CAP / "cap_process_engine.png", "Process Engine"),
    ]
    cw, ch = 3.95, 2.45
    x0, y0, gx, gy = 0.62, 1.35, 0.2, 0.35
    for i, (p, lab) in enumerate(grid):
        r, c = divmod(i, 3)
        x = x0 + c * (cw + gx); y = y0 + r * (ch + gy + 0.15)
        rect(s, x, y, cw, ch, WHITE, line=RGBColor(0xDD, 0xE3, 0xE8), lw=1, rounded=True)
        image_fit(s, p, x + 0.08, y + 0.08, cw - 0.16, ch - 0.16)
        text(s, x, y + ch + 0.02, cw, 0.3,
             [{"text": lab, "size": 11.5, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}])
    footer(s)
    notes(s, "L'application comporte six pages : Supervision (l'accueil-vitrine), Profile (le profil de "
              "vis et les zones thermiques), Settings (paramètres procédé et persistance), Run Analysis "
              "(analyse d'un essai), History (historique), et Process Engine (vue moteur en lecture "
              "seule). Charte sombre, industrielle ; chaque bloc pilote le moteur — rien de décoratif.")
    SCRIPT.append(("Slide 10 — Interface développée (≈ 40 s)",
        "Voici l'application réellement développée et déployée, en six pages. Supervision est l'accueil : "
        "statut machine, score de stabilité, alertes et recommandations. Profile sert à configurer le "
        "profil de vis et les zones thermiques. Settings regroupe les paramètres procédé et la "
        "persistance. Run Analysis analyse un essai dans le temps. History conserve l'historique des "
        "configurations. Et Process Engine offre la vue moteur en lecture seule. La charte est sombre et "
        "industrielle, et — exigence forte du projet — chaque bloc de l'interface pilote le moteur : "
        "rien n'est purement décoratif."))

    # ---- S11 Résultats & démonstration ---------------------------------- #
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT); top_band(s, "Résultats & démonstration", 11)
    image_fit(s, V2 / "figv2_05_before_after_dashboard.png", 0.55, 1.45, 7.6, 3.6)
    caption(s, 0.55, 5.1, 7.6, "Avant → après : cas C3 (à risque) vs C5 (recommandation appliquée)")
    text(s, 8.4, 1.5, 4.4, 5.0, [
        {"text": "Boucle complète validée", "size": 16, "bold": True, "color": NAVY, "space_after": 8},
        {"runs": [("C1 ", {"bold": True, "color": GREEN}), ("référence lithiée — 65/100", {"color": INK})], "size": 13.5, "space_after": 5},
        {"runs": [("C2 ", {"bold": True, "color": GREEN}), ("optimisée — 82/100", {"color": INK})], "size": 13.5, "space_after": 5},
        {"runs": [("C3 ", {"bold": True, "color": RED}), ("surcharge LATP — 46/100, alerte Z5", {"color": INK})], "size": 13.5, "space_after": 5},
        {"runs": [("C4 ", {"bold": True, "color": AMBER}), ("recommandation chiffrée de l'agent", {"color": INK})], "size": 13.5, "space_after": 5},
        {"runs": [("C5 ", {"bold": True, "color": GREEN}), ("après correction — 78/100, alerte levée", {"color": INK})], "size": 13.5, "space_after": 10},
        {"runs": [("Gain C3 → C5 : ", {"size": 14, "bold": True, "color": NAVY}),
                  ("+32 points, p 0,35 → 0,87", {"size": 14, "bold": True, "color": GREEN})], "line": 1.05},
    ])
    footer(s)
    notes(s, "La démonstration suit cinq cas. C1, formulation lithiée de référence, sert de point zéro. "
              "C2 montre un gain par la géométrie. C3 provoque une surcharge céramique : le score chute à "
              "46, alerte rouge en Z5. C4 propose un plan chiffré. C5 valide la correction : score 78, "
              "alerte levée. Soit +32 points et une probabilité de stabilité qui passe de 0,35 à 0,87. "
              "La boucle détection → recommandation → correction est démontrée de bout en bout.")
    SCRIPT.append(("Slide 11 — Résultats & démonstration (≈ 55 s)",
        "La démonstration repose sur cinq cas d'usage lithiés, figés à l'avance pour être reproductibles. "
        "C1 est la formulation de référence, à 65 sur 100 : mon point zéro. C2, optimisée par la "
        "géométrie de vis, monte à 82. C3 provoque délibérément une surcharge en électrolyte céramique : "
        "le score chute à 46, le régime devient instable et une alerte rouge se déclenche en zone Z5, "
        "étayée par un couple élevé et un remplissage critique. En C4, l'agent propose un plan d'actions "
        "chiffré et hiérarchisé. En C5, après application, le score remonte à 78, l'alerte est levée. "
        "Au total, de C3 à C5 : plus 32 points de score et une probabilité de stabilité qui passe de "
        "0,35 à 0,87. La boucle complète — détecter, recommander, corriger, vérifier — fonctionne."))

    # ---- S12 Limites ---------------------------------------------------- #
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT); top_band(s, "Limites assumées", 12)
    cards = [
        ("Données rares", "8 essais exploitables — robustesse statistique bornée", AMBER),
        ("Valeurs nominales", "Modèle non calibré industriellement (tendance relative)", AMBER),
        ("Équations différées", "E5 / E6 / E7 renvoient None, documentées comme telles", BLUE2),
        ("Mono-opérateur", "Pas de gestion multi-utilisateur concurrente", BLUE2),
        ("Dette d'industrialisation", "Pas de CI/CD, couverture non mesurée formellement", GREY),
        ("Posture", "Aide à la décision — jamais pilotage automatique", GREEN),
    ]
    cw, ch = 3.95, 2.35
    for i, (t, b, ac) in enumerate(cards):
        r, c = divmod(i, 3)
        x = 0.62 + c * (cw + 0.2); y = 1.5 + r * (ch + 0.3)
        card(s, x, y, cw, ch, t, [b], accent=ac, title_size=15, body_size=12.5)
    footer(s)
    notes(s, "Je présente les limites sans détour : peu de données (8 essais), des valeurs nominales non "
              "calibrées à lire en tendance, des équations E5/E6/E7 volontairement différées, un modèle "
              "mono-opérateur, une dette d'industrialisation (pas de CI/CD). Et une posture constante : "
              "c'est une aide à la décision, jamais un pilotage automatique. Cette honnêteté est, "
              "paradoxalement, un facteur de crédibilité.")
    SCRIPT.append(("Slide 12 — Limites assumées (≈ 40 s)",
        "L'honnêteté intellectuelle fait partie de l'évaluation. J'assume donc clairement les limites. "
        "D'abord, la rareté des données : huit essais seulement. Ensuite, des valeurs nominales, non "
        "calibrées industriellement, à lire en tendance relative. Plusieurs équations — énergie "
        "mécanique spécifique locale, température réelle avancée, pression filière — sont volontairement "
        "différées et renvoient explicitement une valeur nulle. Le modèle est mono-opérateur, et il "
        "reste une dette d'industrialisation, sans intégration continue. Enfin, une posture constante : "
        "l'outil est une aide à la décision, jamais un pilotage automatique. La décision reste humaine."))

    # ---- S13 Apports pour Rondol ---------------------------------------- #
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT); top_band(s, "Apports pour Rondol Industrie", 13)
    cards = [
        ("Différenciation", "Une aide à la décision explicable, rare sur le segment SSB", NAVY),
        ("Aide opérateur", "Lecture claire du procédé, anticipation des dérives", BLUE2),
        ("Réduction du risque essai", "Comparer des configurations avant de consommer matière", GREEN),
        ("Valorisation de l'IA", "Transfert du savoir-faire pharma (HME) vers l'énergie", GREEN),
    ]
    cw, ch = 5.95, 2.35
    for i, (t, b, ac) in enumerate(cards):
        r, c = divmod(i, 2)
        x = 0.62 + c * (cw + 0.25); y = 1.55 + r * (ch + 0.35)
        card(s, x, y, cw, ch, t, [b], accent=ac, title_size=17, body_size=14)
    footer(s)
    notes(s, "Pour Rondol, quatre apports : une différenciation concurrentielle (peu d'acteurs offrent "
              "une IA explicable sur le segment batteries), une aide concrète à l'opérateur, une "
              "réduction du risque et du coût d'essai en comparant les configurations en amont, et la "
              "valorisation de l'IA comme prolongement du savoir-faire pharmaceutique vers l'énergie.")
    SCRIPT.append(("Slide 13 — Apports pour Rondol (≈ 35 s)",
        "Sur le plan stratégique, l'outil apporte quatre choses à Rondol. Une différenciation : très peu "
        "de concurrents proposent une aide à la décision explicable sur le segment des batteries "
        "tout-solide. Une aide concrète à l'opérateur, qui lit l'état du procédé et anticipe les "
        "dérives. Une réduction du risque et du coût d'essai, en comparant des configurations avant de "
        "consommer de la matière coûteuse. Et la valorisation de l'IA comme transfert du savoir-faire "
        "pharmaceutique de Rondol vers les matériaux d'énergie."))

    # ---- S14 Perspectives ----------------------------------------------- #
    s = blank(prs); rect(s, 0, 0, SW, SH, LIGHT); top_band(s, "Perspectives", 14)
    image_fit(s, FIG / "fig_roadmap.png", 0.7, 1.5, 11.9, 4.4)
    caption(s, 0.7, 5.95, 11.9, "Feuille de route en quatre axes")
    footer(s)
    notes(s, "Quatre axes d'évolution : enrichir les données et engager une calibration industrielle ; "
              "coder les équations différées et ajouter l'interprétabilité SHAP ; résorber la dette "
              "logicielle avec une CI/CD ; ouvrir au multi-utilisateur et intégrer des capteurs temps "
              "réel de couple et de pression pour passer des proxys V1 à des mesures directes V2.")
    SCRIPT.append(("Slide 14 — Perspectives (≈ 35 s)",
        "Les limites dessinent une feuille de route en quatre axes. Premier axe : enrichir le jeu de "
        "données et engager une vraie calibration industrielle. Deuxième : matérialiser les équations "
        "aujourd'hui différées et ajouter l'interprétabilité SHAP au modèle. Troisième : résorber la "
        "dette logicielle avec une chaîne d'intégration continue. Quatrième : ouvrir au multi-utilisateur "
        "et intégrer des capteurs temps réel de couple et de pression, pour passer des proxys de la V1 à "
        "des mesures directes en V2 — la condition d'une crédibilité pleinement industrielle."))

    # ---- S15 Conclusion (fond bleu) ------------------------------------- #
    s = blank(prs); rect(s, 0, 0, SW, SH, NAVY); rect(s, 0, 0, 0.28, SH, GREEN2)
    text(s, 1.0, 1.2, SW - 2.0, 0.6,
         [{"text": "Conclusion", "size": 26, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}])
    rect(s, 1.6, 2.1, SW - 3.2, 1.7, NAVY_D, line=RGBColor(0x9D, 0xC3, 0xE6), lw=1.2, rounded=True)
    text(s, 2.1, 2.3, SW - 4.2, 1.3,
         [{"text": "Un prototype professionnel d'aide à la décision — pas un outil magique, "
                   "mais une base solide, honnête et démontrable.",
           "size": 19, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "line": 1.15}],
         anchor=MSO_ANCHOR.MIDDLE)
    for i, (t, b) in enumerate([
        ("Données réelles", "valorisées en aide à la décision"),
        ("Deux niveaux d'IA", "assumés et argumentés"),
        ("Honnêteté", "méthodologique de bout en bout")]):
        x = 1.1 + i * 3.8
        text(s, x, 4.2, 3.6, 1.3, [
            {"text": t, "size": 15, "bold": True, "color": GREEN2, "align": PP_ALIGN.CENTER, "space_after": 3},
            {"text": b, "size": 12.5, "color": RGBColor(0xCF, 0xDD, 0xEC), "align": PP_ALIGN.CENTER, "line": 1.0}])
    text(s, 1.0, 5.7, SW - 2.0, 0.6,
         [{"text": "Merci de votre attention.", "size": 18, "italic": True, "color": WHITE, "align": PP_ALIGN.CENTER}])
    logo_chip(s, NEXA, SW / 2 - 1.7, 6.45, 0.55)
    iw, ih = Image.open(RONDOL).size
    logo_chip(s, RONDOL, SW / 2 + 0.4, 6.5, 0.48)
    notes(s, "En conclusion : oui, il est possible de transformer des données d'instrumentation réelles "
              "en un outil cohérent qui rend lisible le procédé, anticipe les dérives et recommande des "
              "actions justifiées. Ce n'est pas un outil magique, mais un prototype professionnel d'aide "
              "à la décision — solide, honnête et démontrable. Merci ; je suis à votre disposition pour "
              "vos questions.")
    SCRIPT.append(("Slide 15 — Conclusion (≈ 35 s)",
        "Pour conclure, la réponse à ma problématique est positive et nuancée : il est possible de "
        "transformer un gisement de données réelles, jusque-là sous-exploité, en un outil cohérent qui "
        "rend lisible l'état du procédé, anticipe ses dérives et recommande des actions chiffrées et "
        "justifiées. Ce n'est pas un outil magique : c'est un prototype professionnel d'aide à la "
        "décision — une base solide, honnête et démontrable, et un vrai facteur de différenciation pour "
        "Rondol Industrie. Je vous remercie de votre attention et reste à votre disposition pour vos "
        "questions."))

    return prs


# --------------------------------------------------------------------------- #
def build_script_docx():
    from docx import Document
    from docx.shared import Pt as DPt, RGBColor as DRGB, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc = Document()
    n = doc.styles["Normal"]; n.font.name = "Calibri"; n.font.size = DPt(11.5)
    h = doc.add_paragraph()
    r = h.add_run("Script de soutenance — Mémoire de thèse professionnelle")
    r.bold = True; r.font.size = DPt(16); r.font.color.rgb = DRGB(0x1F, 0x4E, 0x79)
    sub = doc.add_paragraph()
    rs = sub.add_run("Wilfried Galtier MBEUMI — Nexa Digital School × Rondol Industrie — durée cible ≈ 10 minutes")
    rs.italic = True; rs.font.size = DPt(11); rs.font.color.rgb = DRGB(0x59, 0x59, 0x59)
    doc.add_paragraph()
    for titre, discours in SCRIPT:
        p = doc.add_paragraph()
        rt = p.add_run(titre); rt.bold = True; rt.font.size = DPt(13); rt.font.color.rgb = DRGB(0x2E, 0x75, 0xB6)
        p.paragraph_format.space_before = DPt(10)
        body = doc.add_paragraph(discours)
        body.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        body.paragraph_format.line_spacing = 1.25
    tip = doc.add_paragraph()
    rtip = tip.add_run("Conseil : ~40 secondes par slide en moyenne ; ralentir sur les slides 8 (ML) "
                       "et 11 (résultats), accélérer sur 4, 13 et 14.")
    rtip.italic = True; rtip.font.size = DPt(10.5); rtip.font.color.rgb = DRGB(0x59, 0x59, 0x59)
    tip.paragraph_format.space_before = DPt(14)
    doc.save(str(SCRIPT_DOCX))


def main():
    prs = build(prs_new())
    prs.save(str(PPTX))
    print(f"[OK] PPTX : {PPTX}  ({len(prs.slides)} slides)")
    build_script_docx()
    print(f"[OK] Script oral : {SCRIPT_DOCX}")


if __name__ == "__main__":
    main()
