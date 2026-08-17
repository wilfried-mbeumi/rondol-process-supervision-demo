# -*- coding: utf-8 -*-
"""Soutenance PARTIE 2 (diapos 11–20) — design professionnel.

Usage : python scripts/build_soutenance_part2.py
Fusion : ouvrir les deux PPTX et copier les diapos de la partie 2 après la 10.
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "memoire_figures"
CAP = ROOT / "figures_memoire"
OUT = ROOT / "reports" / "soutenance" / "PREZ_part2.pptx"

# ── Palette (identique part 1) ────────────────────────────────────────────── #
BLEU       = RGBColor(0x1B, 0x3A, 0x5C)
BLEU_MED   = RGBColor(0x2E, 0x75, 0xB6)
BLEU_CLAIR = RGBColor(0xD6, 0xE4, 0xF0)
VERT       = RGBColor(0x1B, 0x7A, 0x3D)
VERT_CLAIR = RGBColor(0xD5, 0xF0, 0xDF)
ROUGE      = RGBColor(0xB0, 0x3A, 0x2E)
ROUGE_CLAIR= RGBColor(0xF9, 0xE2, 0xDF)
AMBRE      = RGBColor(0xB7, 0x79, 0x1F)
AMBRE_CLAIR= RGBColor(0xFD, 0xF2, 0xDB)
GRIS       = RGBColor(0x6B, 0x6B, 0x6B)
GRIS_CLAIR = RGBColor(0xF2, 0xF2, 0xF2)
ENCRE      = RGBColor(0x1E, 0x1E, 0x1E)
BLANC      = RGBColor(0xFF, 0xFF, 0xFF)
FOND       = RGBColor(0xFD, 0xFD, 0xFD)

POLICE = "Calibri"
L, H = Inches(13.333), Inches(7.5)
MARGE = Inches(0.75)
MARGE_D = Inches(0.75)
UTILE = L - MARGE - MARGE_D

# ── Helpers (même API que part 1) ─────────────────────────────────────────── #

def _rect(d, x, y, w, h, fill=None, border=None, border_w=Pt(1), radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = d.shapes.add_shape(shape_type, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = border_w
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def _txt(d, x, y, w, h, texte, taille=16, gras=False, couleur=ENCRE,
         align=PP_ALIGN.LEFT, italique=False, interligne=1.3, ancre=None):
    zone = d.shapes.add_textbox(x, y, w, h)
    tf = zone.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = 0
    if ancre:
        tf.vertical_anchor = ancre
    p = tf.paragraphs[0]
    p.alignment = align; p.line_spacing = interligne
    r = p.add_run()
    r.text = texte; r.font.size = Pt(taille); r.font.bold = gras
    r.font.italic = italique; r.font.color.rgb = couleur; r.font.name = POLICE
    return zone, tf


def _add_para(tf, texte, taille=16, gras=False, couleur=ENCRE, avant=6,
              italique=False, interligne=1.3):
    p = tf.add_paragraph()
    p.space_before = Pt(avant); p.line_spacing = interligne
    r = p.add_run()
    r.text = texte; r.font.size = Pt(taille); r.font.bold = gras
    r.font.italic = italique; r.font.color.rgb = couleur; r.font.name = POLICE
    return p


def slide(prs):
    d = prs.slides.add_slide(prs.slide_layouts[6])
    d.background.fill.solid(); d.background.fill.fore_color.rgb = FOND
    return d


def header(d, titre, sous_titre=None, accent=BLEU):
    _rect(d, Inches(0), Inches(0), L, Inches(1.35), fill=accent)
    _rect(d, Inches(0), Inches(1.35), L, Inches(0.06), fill=VERT)
    _txt(d, MARGE, Inches(0.22), UTILE, Inches(0.85), titre,
         taille=28, gras=True, couleur=BLANC, interligne=1.05)
    if sous_titre:
        _txt(d, MARGE, Inches(1.5), UTILE, Inches(0.4), sous_titre,
             taille=13, couleur=GRIS, italique=True)


def notes(d, txt):
    d.notes_slide.notes_text_frame.text = txt.strip()


def kpi_card(d, x, y, w, h, valeur, label, couleur=BLEU, fond=BLANC, border=None):
    _rect(d, x, y, w, h, fill=fond, border=border or couleur, border_w=Pt(1.5), radius=True)
    _txt(d, x, y + Inches(0.18), w, Inches(0.7), valeur,
         taille=44, gras=True, couleur=couleur, align=PP_ALIGN.CENTER, interligne=1.0)
    _txt(d, x + Inches(0.1), y + Inches(0.95), w - Inches(0.2), Inches(0.9), label,
         taille=12, couleur=GRIS, align=PP_ALIGN.CENTER, interligne=1.25)


def kpi_row(d, items, y=Inches(2.3), couleur=BLEU, h=Inches(1.7)):
    n = len(items)
    gap = Inches(0.3)
    w = Emu(int((UTILE - gap * (n - 1)) / n))
    for i, (val, lab) in enumerate(items):
        x = Emu(int(MARGE + i * (w + gap)))
        kpi_card(d, x, y, w, h, val, lab, couleur=couleur)


def content_card(d, x, y, w, h, fill=BLANC, border=None):
    _rect(d, x, y, w, h, fill=fill, border=border, border_w=Pt(1), radius=True)


def blockquote(d, texte, y=Inches(5.5), couleur=BLEU, w=None):
    w = w or UTILE
    _rect(d, MARGE, y, Inches(0.07), Inches(1.2), fill=couleur)
    bg = BLEU_CLAIR if couleur == BLEU else VERT_CLAIR if couleur == VERT else \
         ROUGE_CLAIR if couleur == ROUGE else AMBRE_CLAIR
    content_card(d, MARGE, y, w, Inches(1.2), fill=bg, border=couleur)
    _txt(d, MARGE + Inches(0.3), y + Inches(0.15), w - Inches(0.5), Inches(0.9),
         texte, taille=15, gras=True, couleur=couleur, interligne=1.3)


def bullet_block(d, items, y=Inches(2.2), taille=18, w=None, x=None, accent=BLEU):
    x = x or MARGE
    w = w or UTILE
    cur_y = y
    for item in items:
        if isinstance(item, tuple):
            titre, corps = item
        else:
            titre, corps = None, item
        card_h = Inches(0.85) if len(corps) < 120 else Inches(1.1)
        content_card(d, x, cur_y, w, card_h, fill=BLANC, border=GRIS_CLAIR)
        _rect(d, x, cur_y, Inches(0.06), card_h, fill=accent)
        text_x = x + Inches(0.22)
        text_w = w - Inches(0.35)
        if titre:
            _, tf = _txt(d, text_x, cur_y + Inches(0.1), text_w, Inches(0.35),
                         titre, taille=taille, gras=True, couleur=accent)
            _add_para(tf, corps, taille=taille - 2, couleur=ENCRE, avant=2)
        else:
            _txt(d, text_x, cur_y + Inches(0.12), text_w, card_h - Inches(0.2),
                 corps, taille=taille - 1, couleur=ENCRE)
        cur_y += card_h + Inches(0.12)


def fig(d, chemin, haut=Inches(2.0), bas=Inches(0.3)):
    from PIL import Image as PILImage
    iw, ih = PILImage.open(chemin).size
    dispo_h = H - haut - bas
    dispo_w = UTILE
    ratio = iw / ih
    w = dispo_w; h = Emu(int(w / ratio))
    if h > dispo_h:
        h = dispo_h; w = Emu(int(h * ratio))
    content_card(d, Emu(int((L - w) / 2)) - Inches(0.1), haut - Inches(0.08),
                 Emu(int(w)) + Inches(0.2), Emu(int(h)) + Inches(0.16),
                 fill=BLANC, border=GRIS_CLAIR)
    d.shapes.add_picture(str(chemin), Emu(int((L - w) / 2)), haut,
                         width=Emu(int(w)), height=Emu(int(h)))


def num_slide(d, numero, total=20):
    _txt(d, L - Inches(1.0), H - Inches(0.4), Inches(0.7), Inches(0.3),
         f"{numero}/{total}", taille=10, couleur=GRIS, align=PP_ALIGN.RIGHT)


# =========================================================================== #
prs = Presentation()
prs.slide_width, prs.slide_height = L, H

# ── 11. PREMIER RÉSULTAT ─────────────────────────────────────────────────── #
d = slide(prs)
header(d, "Premier résultat : ce que je croyais avoir obtenu",
       "Augmentation de données — 800 fenêtres synthétiques depuis l'échantillon réel",
       accent=VERT)
num_slide(d, 11)

bullet_block(d, [
    ("Le constat.",
     "Huit essais, c'est trop peu. J'ai généré 800 fenêtres synthétiques par bootstrap "
     "conditionné par classe, avec jitter borné et reproduction des imperfections capteur "
     "— injectées à l'entraînement uniquement."),
], y=Inches(2.1), accent=VERT)

kpi_row(d, [("0,809 → 0,918", "F1-macro\ndu Random Forest"),
            ("÷ 3", "écart-type inter-essais\n(0,176 → 0,054)")],
        y=Inches(3.6), couleur=VERT, h=Inches(1.6))

blockquote(d, "J'ai présenté ce résultat comme l'aboutissement du volet prédictif. "
              "Il était faux.", y=Inches(5.7), couleur=ROUGE)
notes(d, """Face au manque de données, j'ai construit une augmentation à partir de
l'échantillon réel. Le gain paraissait spectaculaire : le Random Forest passait de
0,809 à 0,918, et sa variance était divisée par plus de trois.
MARQUEZ UN TEMPS D'ARRÊT ICI. Puis : « Il était faux. »
[~2 minutes]""")

# ── 12. L'AUDIT ──────────────────────────────────────────────────────────── #
d = slide(prs)
d.background.fill.solid(); d.background.fill.fore_color.rgb = RGBColor(0xFD, 0xF6, 0xF5)
header(d, "L'audit : une fuite par ancrage", accent=ROUGE)
num_slide(d, 12)

bullet_block(d, [
    ("Le défaut.",
     "Le pool synthétique était généré une seule fois, à partir des huit "
     "essais réels — puis réutilisé dans chaque pli."),
    ("La conséquence.",
     "L'essai censé être exclu contribuait indirectement à l'entraînement : "
     "ses fenêtres servaient de points d'ancrage au bootstrap."),
    ("Pourquoi c'était invisible.",
     "Le pli de test ne contenait aucune fenêtre synthétique. À la lecture du code, tout semblait correct."),
    ("La correction.",
     "Régénérer le pool à chaque pli, à partir des seuls essais d'entraînement."),
], y=Inches(2.1), taille=18, accent=ROUGE)

notes(d, """Voici ce que l'audit a révélé.
Le pool synthétique était généré UNE FOIS à partir des huit essais, puis réutilisé
dans chaque pli. Donc l'essai que je gardais pour tester avait servi à fabriquer les
données d'entraînement.
C'est ce qu'on appelle une fuite par ancrage.
La correction tient en une phrase : régénérer le pool à chaque pli.
[~3 minutes — c'est le cœur de la soutenance, prenez le temps]""")

# ── 13. LE VRAI RÉSULTAT ─────────────────────────────────────────────────── #
d = slide(prs)
header(d, "Le vrai résultat : le gain disparaît",
       "F1-macro LOGO — sans augmentation · globale (fuitée) · par pli (corrigée)",
       accent=ROUGE)
num_slide(d, 13)

lignes = [("SVM (RBF)",               "0,805", "0,868", "0,824"),
          ("Régression logistique",    "0,799", "0,860", "0,809"),
          ("Random Forest",            "0,809", "0,918", "0,809"),
          ("XGBoost",                  "0,757", "0,900", "0,801"),
          ("Réseau de neurones (MLP)", "0,778", "0,862", "0,781")]

y0 = Inches(2.15)
row_h = Inches(0.52)
col_x = [MARGE + Inches(0.15), MARGE + Inches(4.5), MARGE + Inches(7.0), MARGE + Inches(9.5)]
col_w = [Inches(4.2), Inches(2.3), Inches(2.3), Inches(2.3)]

content_card(d, MARGE, y0 - Inches(0.1), UTILE, Inches(0.45) + row_h * 5 + Inches(0.2),
             fill=BLANC, border=GRIS_CLAIR)

_rect(d, MARGE, y0 - Inches(0.1), UTILE, Inches(0.45), fill=BLEU)
for j, lib in enumerate(["Modèle", "Sans augm.", "Globale (fuitée)", "Par pli (corrigée)"]):
    _txt(d, col_x[j], y0 - Inches(0.06), col_w[j], Inches(0.35), lib,
         taille=12, gras=True, couleur=BLANC)

for i, (m, a, b, c) in enumerate(lignes):
    y = y0 + Inches(0.45) + i * row_h
    fort = (m == "Random Forest")
    if fort:
        _rect(d, MARGE + Inches(0.05), y - Inches(0.03), UTILE - Inches(0.1),
              row_h - Inches(0.02), fill=BLEU_CLAIR)
    if i % 2 == 1 and not fort:
        _rect(d, MARGE + Inches(0.05), y - Inches(0.03), UTILE - Inches(0.1),
              row_h - Inches(0.02), fill=GRIS_CLAIR)
    _txt(d, col_x[0], y, col_w[0], Inches(0.4), m, taille=14, gras=fort,
         couleur=BLEU if fort else ENCRE)
    _txt(d, col_x[1], y, col_w[1], Inches(0.4), a, taille=14, couleur=GRIS)
    _txt(d, col_x[2], y, col_w[2], Inches(0.4), b, taille=14, gras=fort,
         couleur=ROUGE if fort else GRIS)
    _txt(d, col_x[3], y, col_w[3], Inches(0.4), c, taille=14, gras=True,
         couleur=VERT)

content_card(d, MARGE, Inches(5.3), UTILE, Inches(1.5), fill=GRIS_CLAIR)
_txt(d, MARGE + Inches(0.2), Inches(5.4), UTILE - Inches(0.4), Inches(1.3),
     "Sur le Random Forest, le gain passe de +0,109 à −0,001. Aucun modèle "
     "n'atteint 0,85 ; les cinq restent groupés entre 0,78 et 0,82, pour des "
     "écarts-types de 0,13 à 0,17 : statistiquement indiscernables.",
     taille=15, couleur=ENCRE, interligne=1.35)
notes(d, """Regardez la ligne Random Forest : 0,809 sans augmentation, 0,918 avec
l'augmentation fuitée, et de nouveau 0,809 une fois la fuite corrigée.
Le gain passe de +0,109 à −0,001. Il disparaît.
Aucun modèle n'atteint 0,85. Les cinq sont statistiquement indiscernables.
[~3 minutes]""")

# ── 14. POURQUOI JE LE DIS ───────────────────────────────────────────────── #
d = slide(prs)
header(d, "Pourquoi je l'expose plutôt que de le taire", accent=ROUGE)
num_slide(d, 14)

bullet_block(d, [
    ("Déontologique.",
     "Annoncer 0,918 à Rondol aurait été une promesse que le modèle "
     "n'aurait pas tenue en production."),
    ("Méthodologique.",
     "Cette fuite illustre mieux que n'importe quel développement réussi "
     "le point central de mon travail : la performance est une propriété du "
     "protocole d'évaluation autant que de l'algorithme."),
    ("Pratique.",
     "Prédictions par pli, métriques et générateur corrigé sont livrés "
     "dans le dépôt. Chaque chiffre est recalculable et contestable."),
], y=Inches(2.1), taille=19, accent=ROUGE)

blockquote(d, "Le Random Forest reste le modèle retenu — non pour son score, mais pour "
              "son interprétabilité, sa tolérance aux capteurs manquants et sa stabilité "
              "inter-essais.", y=Inches(5.6), couleur=BLEU)
notes(d, """Trois raisons de dire les choses plutôt que de les taire.
Déontologique : Rondol aurait construit des décisions sur un chiffre qui ne tenait pas.
Méthodologique : la performance d'un modèle est une propriété du protocole autant que de l'algorithme.
Pratique : tout est livré, donc vérifiable.
[~2 minutes]""")

# ── 15. VALIDATION EXTERNE ───────────────────────────────────────────────── #
d = slide(prs)
header(d, "Épreuve de transférabilité : validation externe",
       "Base continue de 100 800 lignes — modèle appliqué sans réentraînement")
num_slide(d, 15)

kpi_row(d, [("3 479", "fenêtres évaluées\nsur 15 runs simulés"),
            ("0,753", "AUC — le pouvoir\ndiscriminant subsiste"),
            ("62 %", "des fenêtres instables\ndétectées")],
        y=Inches(2.2), h=Inches(1.65))

content_card(d, MARGE, Inches(4.4), UTILE, Inches(2.3), fill=GRIS_CLAIR)
_txt(d, MARGE + Inches(0.2), Inches(4.5), UTILE - Inches(0.4), Inches(2.1),
     "Les erreurs sont majoritairement conservatrices : fausses alertes plutôt que "
     "dérives manquées.\n\n"
     "La génération n'a volontairement pas été ajustée pour flatter ces chiffres — "
     "l'écart avec les essais réels mesure la sensibilité au changement de "
     "distribution, et confirme le statut : indicateur d'aide à la décision, "
     "pas détecteur certifié.",
     taille=14, couleur=GRIS, italique=True, interligne=1.35)
notes(d, """Pour tester la généralisation sans attendre une nouvelle campagne, j'ai
généré une base continue et appliqué le modèle sans réentraînement.
L'AUC tombe à 0,753 — le pouvoir discriminant subsiste, mais il baisse.
Le point important : je n'ai pas ajusté la génération pour améliorer ces chiffres.
Les erreurs vont dans le bon sens : plutôt une fausse alerte qu'une dérive manquée.
[~2 minutes]""")

# ── 16. L'APPLICATION ────────────────────────────────────────────────────── #
d = slide(prs)
header(d, "L'application : une HMI industrielle, pas un tableau de bord",
       "7 pages · accès protégé · bilingue FR/EN · persistance en trois couches")
num_slide(d, 16)

fig(d, CAP / "cap_supervision.png", haut=Inches(2.1), bas=Inches(1.1))

content_card(d, MARGE, Inches(6.6), UTILE, Inches(0.55), fill=GRIS_CLAIR)
_txt(d, MARGE + Inches(0.15), Inches(6.65), UTILE - Inches(0.3), Inches(0.4),
     "rondol-process-supervision-demo.streamlit.app   ·   "
     "contrastes WCAG 2.1 AA vérifiés de 6,38:1 à 18,39:1",
     taille=11, couleur=GRIS, align=PP_ALIGN.CENTER)
notes(d, """Voici la page Supervision, celle que l'opérateur ouvre en premier.
Score de stabilité, probabilité de dérive, alertes et recommandations — tout est sur un seul écran.
L'application est en ligne et accessible publiquement.
[~2 minutes]""")

# ── 17. LES CINQ CAS ────────────────────────────────────────────────────── #
d = slide(prs)
header(d, "Cinq cas pour prouver que l'outil réagit juste",
       "Sensibilité · détectabilité · réversibilité")
num_slide(d, 17)

cas = [("C1", "Formulation lithiée\nde référence", "65/100", BLEU,
        "LFP 65 % · PVDF 8 %\nLATP 17 %"),
       ("C2", "Configuration\noptimisée", "82/100", VERT, ""),
       ("C3", "Défaut\nprovoqué", "46/100", ROUGE, "Alerte rouge\nlocalisée en Z5"),
       ("C4", "Recommandation\nde l'agent", "→", AMBRE, "Quel paramètre,\nde combien"),
       ("C5", "Après\ncorrection", "+32 pts", VERT, "Alerte levée\n+0,52 proba")]
n = len(cas)
gap = Inches(0.2)
w = Emu(int((UTILE - gap * (n - 1)) / n))
for i, (code, desc, val, coul, detail) in enumerate(cas):
    x = Emu(int(MARGE + i * (w + gap)))
    content_card(d, x, Inches(2.15), w, Inches(3.3), fill=BLANC, border=coul)
    _rect(d, x, Inches(2.15), w, Inches(0.55), fill=coul)
    _txt(d, x, Inches(2.2), w, Inches(0.45), code,
         taille=20, gras=True, couleur=BLANC, align=PP_ALIGN.CENTER)
    _txt(d, x + Inches(0.08), Inches(2.85), w - Inches(0.16), Inches(0.7), desc,
         taille=12, couleur=ENCRE, align=PP_ALIGN.CENTER, interligne=1.25)
    _txt(d, x, Inches(3.55), w, Inches(0.55), val,
         taille=28, gras=True, couleur=coul, align=PP_ALIGN.CENTER)
    if detail:
        _txt(d, x + Inches(0.08), Inches(4.15), w - Inches(0.16), Inches(0.8), detail,
             taille=10, couleur=GRIS, align=PP_ALIGN.CENTER, interligne=1.25)

blockquote(d, "La transition C3 → C5 : +32 points de score, +0,52 de probabilité de "
              "stabilité, alerte levée — et la projection annoncée par l'agent en C4 "
              "se vérifie.", y=Inches(5.8), couleur=VERT)
notes(d, """Cinq cas construits pour démontrer trois propriétés.
Sensibilité : de C1 à C2, l'outil réagit — 65 puis 82 sur 100.
Détectabilité : en C3, défaut provoqué, score chute à 46, alerte rouge en Z5.
Réversibilité : en C4 recommandation chiffrée, en C5 score remonte de 32 points.
[~2 minutes 30]""")

# ── 18. PILOTAGE ─────────────────────────────────────────────────────────── #
d = slide(prs)
header(d, "Piloter le projet comme un chef de projet",
       "Kanban à encours limité · jalons tenus · risques cartographiés")
num_slide(d, 18)

fig(d, FIG / "fig_gantt.png", haut=Inches(2.1), bas=Inches(1.5))

bullets_pm = [("WIP = 1.", "Aucune brique entamée avant validation de la précédente."),
              ("Jalons tenus.", "Campagne 7–13 avril 2026, démo client 16 juin 2026."),
              ("12 incidents", "tracés, résolus et figés en tests de non-régression.")]
cur_y = Inches(5.8)
gap_x = Inches(0.3)
bw = Emu(int((UTILE - gap_x * 2) / 3))
for i, (titre, desc) in enumerate(bullets_pm):
    bx = Emu(int(MARGE + i * (bw + gap_x)))
    content_card(d, bx, cur_y, bw, Inches(0.95), fill=BLANC, border=BLEU)
    _rect(d, bx, cur_y, bw, Inches(0.04), fill=BLEU)
    _, tf = _txt(d, bx + Inches(0.1), cur_y + Inches(0.12), bw - Inches(0.2), Inches(0.3),
                 titre, taille=13, gras=True, couleur=BLEU)
    _add_para(tf, desc, taille=11, couleur=ENCRE, avant=2, interligne=1.2)

notes(d, """Méthode Kanban à encours limité : une seule brique en cours à la fois.
Les deux jalons datés ont été tenus.
Douze incidents de production tracés avec leur cause, leur correctif et leur commit.
[~2 minutes]""")

# ── 19. QUALITÉ ──────────────────────────────────────────────────────────── #
d = slide(prs)
header(d, "Un prototype vérifié comme un produit",
       "La qualité logicielle au service de la crédibilité scientifique")
num_slide(d, 19)

kpi_row(d, [("720", "tests automatisés\n75 fichiers · 100 % au vert"),
            ("×1030", "accélération mesurée\ntable indexée vs non indexée"),
            ("12", "incidents de production\nchacun devenu un test")],
        y=Inches(2.1), couleur=BLEU, h=Inches(1.65))

fig(d, FIG / "fig_tests.png", haut=Inches(4.2), bas=Inches(0.9))

content_card(d, MARGE, Inches(6.75), UTILE, Inches(0.5), fill=BLEU_CLAIR, border=BLEU)
_txt(d, MARGE + Inches(0.15), Inches(6.8), UTILE - Inches(0.3), Inches(0.4),
     "La qualité logicielle n'est pas un supplément d'âme : "
     "c'est ce qui rend le résultat scientifique crédible.",
     taille=13, gras=True, couleur=BLEU, align=PP_ALIGN.CENTER)
notes(d, """720 tests automatisés répartis en six familles, tous au vert.
La famille accessibilité est récente : elle verrouille le nom accessible du schéma
de vis et les quatre rapports de contraste.
Chaque incident corrigé est devenu un test — il ne peut pas revenir.
[~1 minute 30]""")

# ── 20. CONCLUSION ──────────────────────────────────────────────────────── #
d = slide(prs)
d.background.fill.solid(); d.background.fill.fore_color.rgb = BLEU
_rect(d, Inches(0), Inches(0), Inches(0.35), H, fill=VERT)
_rect(d, Inches(0), H - Inches(0.08), L, Inches(0.08), fill=VERT)
num_slide(d, 20)

_txt(d, Inches(1.2), Inches(0.6), Inches(10.8), Inches(0.6),
     "Ce que je retiens", taille=34, gras=True, couleur=BLANC)

_rect(d, Inches(1.2), Inches(1.35), Inches(2.5), Inches(0.04), fill=VERT)

blocs = [
    ("Ce qui est acquis", VERT,
     "Un jumeau numérique ancré dans la géométrie réelle de la machine, "
     "un agent explicable qui recommande sans décider, un modèle validé "
     "sous protocole strict, un outil démontrable devant le client."),
    ("Ce qui est assumé", AMBRE,
     "Physique nominale non calibrée industriellement, huit essais d'une "
     "seule campagne, pression filière non modélisée."),
    ("Ce que j'ai vraiment appris", ROUGE,
     "J'ai cru avoir débloqué la situation par l'augmentation de données. "
     "C'est en auditant mon propre protocole que j'ai découvert que le gain "
     "était un artefact. J'aurais pu livrer 0,918 : personne ne l'aurait vu."),
]
cur_y = Inches(1.65)
for titre, coul, corps in blocs:
    content_card(d, Inches(1.2), cur_y, Inches(10.8), Inches(1.1),
                 fill=RGBColor(0x24, 0x46, 0x6C), border=coul)
    _rect(d, Inches(1.2), cur_y, Inches(0.07), Inches(1.1), fill=coul)
    _, tf = _txt(d, Inches(1.5), cur_y + Inches(0.1), Inches(10.3), Inches(0.3),
                 titre, taille=15, gras=True, couleur=coul)
    _add_para(tf, corps, taille=14, couleur=BLANC, avant=4, interligne=1.25)
    cur_y += Inches(1.22)

_rect(d, Inches(1.2), Inches(5.45), Inches(0.07), Inches(1.3), fill=VERT)
content_card(d, Inches(1.2), Inches(5.45), Inches(10.8), Inches(1.3),
             fill=RGBColor(0x14, 0x30, 0x4D), border=VERT)
_txt(d, Inches(1.5), Inches(5.55), Inches(10.3), Inches(1.1),
     "Le facteur limitant n'est ni l'algorithme ni la méthode : c'est le nombre "
     "d'essais. Passer d'un indicateur expérimental à un prédicteur industriel "
     "exigera de nouvelles campagnes — pas un modèle plus sophistiqué.",
     taille=17, gras=True, couleur=BLANC, interligne=1.3)

notes(d, """Ma conclusion en trois temps.
Ce qui est acquis, ce qui est assumé — et surtout ce que j'ai appris.
J'ai cru avoir résolu le problème du manque de données. En auditant mon propre
protocole, j'ai découvert que le gain était un artefact.
La phrase finale : le facteur limitant n'est ni l'algorithme ni la méthode,
c'est le nombre d'essais.
Je vous remercie, et je suis à votre disposition pour vos questions.
[~2 minutes]""")

# ── Sauvegarde ────────────────────────────────────────────────────────────── #
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"[OK] {OUT}  —  {sum(1 for _ in prs.slides)} diapositive(s)")
