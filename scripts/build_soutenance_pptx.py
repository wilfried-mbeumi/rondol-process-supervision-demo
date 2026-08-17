# -*- coding: utf-8 -*-
"""Construit le support final de soutenance Rondol (20 diapositives, 16:9).

Le support applique un système visuel sombre et industriel, utilise uniquement
les chiffres vérifiés du mémoire et embarque des notes d'orateur. Aucun visuel
industriel n'est inventé : les captures proviennent de l'application et les
schémas sont redessinés à partir de l'architecture documentée du projet.
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "memoire_figures"
CAP = ROOT / "figures_memoire"
ASSETS = ROOT / "assets"
OUT = ROOT / "reports" / "soutenance" / "MBEUMI_Wilfried_SOUTENANCE_FINAL.pptx"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.72)
FONT = "Aptos"
FONT_DISPLAY = "Aptos Display"

# Palette sombre, compatible avec l'HMI Rondol.
BG = RGBColor(0x07, 0x11, 0x1F)
BG_2 = RGBColor(0x0B, 0x19, 0x2A)
PANEL = RGBColor(0x10, 0x20, 0x33)
PANEL_2 = RGBColor(0x15, 0x2A, 0x40)
WHITE = RGBColor(0xF7, 0xFA, 0xFC)
MUTED = RGBColor(0x9C, 0xAD, 0xBE)
BLUE = RGBColor(0x42, 0xA5, 0xD5)
CYAN = RGBColor(0x62, 0xD2, 0xE8)
GREEN = RGBColor(0x4C, 0xB9, 0x63)
GREEN_2 = RGBColor(0x86, 0xE0, 0x9A)
RED = RGBColor(0xF0, 0x5D, 0x5E)
AMBER = RGBColor(0xF4, 0xB9, 0x42)
DARK = RGBColor(0x0B, 0x13, 0x1E)
GRID = RGBColor(0x1C, 0x31, 0x47)


def add_text(slide, x, y, w, h, value, *, size=20, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT,
             italic=False, margin=0, lines=1.05):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = lines
    r = p.add_run()
    r.text = value
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def add_rich(slide, x, y, w, h, segments, *, size=20, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, lines=1.1):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = lines
    for value, color, bold in segments:
        r = p.add_run()
        r.text = value
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, *, fill=PANEL, line=None, radius=True,
             transparency=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x, y, w, h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.2)
    shape.shadow.inherit = False
    return shape


def add_glow(slide, x, y, d, color, transparency=82):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.fill.transparency = transparency
    s.line.fill.background()
    return s


def base_slide(prs, number, section, *, accent=BLUE, dark=BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = dark
    add_text(slide, M, Inches(0.25), Inches(4.8), Inches(0.24),
             f"RONDOL  /  {section.upper()}", size=10, color=accent, bold=True)
    add_text(slide, Inches(12.15), Inches(0.22), Inches(0.45), Inches(0.28),
             f"{number:02d}", size=11, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    # Barre de progression discrète.
    add_rect(slide, M, Inches(7.18), Inches(11.88), Inches(0.025),
             fill=GRID, radius=False)
    add_rect(slide, M, Inches(7.18), Inches(11.88 * number / 20), Inches(0.025),
             fill=accent, radius=False)
    return slide


def title(slide, value, kicker=None, *, accent=BLUE, y=0.72, size=30):
    add_text(slide, M, Inches(y), Inches(11.85), Inches(0.86), value,
             size=size, color=WHITE, bold=True, font=FONT_DISPLAY, lines=0.95)
    add_rect(slide, M, Inches(y + 0.88), Inches(0.8), Inches(0.055),
             fill=accent, radius=False)
    if kicker:
        add_text(slide, Inches(1.72), Inches(y + 0.76), Inches(10.4), Inches(0.42),
                 kicker, size=14, color=MUTED, italic=True, valign=MSO_ANCHOR.MIDDLE)


def note(slide, value):
    slide.notes_slide.notes_text_frame.text = value.strip()


def card(slide, x, y, w, h, heading, body, *, accent=BLUE, value=None,
         body_size=17):
    add_rect(slide, x, y, w, h, fill=PANEL, line=GRID)
    add_rect(slide, x, y, Inches(0.07), h, fill=accent, radius=False)
    compact = h <= Inches(1.40) and value is None
    if value:
        add_text(slide, x + Inches(0.3), y + Inches(0.28), w - Inches(0.55),
                 Inches(0.7), value, size=38, color=accent, bold=True,
                 font=FONT_DISPLAY)
        hy = y + Inches(1.06)
    elif compact:
        hy = y + Inches(0.16)
    else:
        hy = y + Inches(0.34)
    if compact:
        add_text(slide, x + Inches(0.26), hy, w - Inches(0.48), Inches(0.28),
                 heading, size=14, color=accent, bold=True)
        add_text(slide, x + Inches(0.26), hy + Inches(0.34), w - Inches(0.48),
                 h - Inches(0.58), body, size=body_size, color=WHITE, lines=1.02)
    else:
        add_text(slide, x + Inches(0.3), hy, w - Inches(0.55), Inches(0.42),
                 heading, size=18, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.3), hy + Inches(0.52), w - Inches(0.55),
                 h - (hy - y) - Inches(0.72), body, size=body_size, color=MUTED,
                 lines=1.18)


def metric(slide, x, y, w, value, label, *, accent=BLUE, sub=None):
    add_rect(slide, x, y, w, Inches(1.62), fill=PANEL, line=GRID)
    add_text(slide, x + Inches(0.18), y + Inches(0.18), w - Inches(0.36),
             Inches(0.62), value, size=34, color=accent, bold=True,
             align=PP_ALIGN.CENTER, font=FONT_DISPLAY)
    add_text(slide, x + Inches(0.18), y + Inches(0.86), w - Inches(0.36),
             Inches(0.35), label, size=15, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, x + Inches(0.18), y + Inches(1.22), w - Inches(0.36),
                 Inches(0.25), sub, size=11, color=MUTED,
                 align=PP_ALIGN.CENTER)


def chip(slide, x, y, w, value, *, fill=PANEL_2, color=CYAN):
    add_rect(slide, x, y, w, Inches(0.42), fill=fill, line=None)
    add_text(slide, x, y + Inches(0.02), w, Inches(0.32), value, size=12,
             color=color, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)


def arrow(slide, x1, y1, x2, y2, *, color=BLUE, width=2.3):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    c.line.end_arrowhead = True
    return c


def picture_contain(slide, path, x, y, w, h, *, background=WHITE, pad=0.14,
                    line=GRID):
    add_rect(slide, x, y, w, h, fill=background, line=line)
    with Image.open(path) as im:
        iw, ih = im.size
    avail_w = w - Inches(2 * pad)
    avail_h = h - Inches(2 * pad)
    ratio = iw / ih
    pw = avail_w
    ph = Emu(int(pw / ratio))
    if ph > avail_h:
        ph = avail_h
        pw = Emu(int(ph * ratio))
    px = x + Emu(int((w - pw) / 2))
    py = y + Emu(int((h - ph) / 2))
    slide.shapes.add_picture(str(path), px, py, width=pw, height=ph)


def three_columns():
    gap = Inches(0.25)
    w = Emu(int((W - 2 * M - 2 * gap) / 3))
    return w, gap


prs = Presentation()
prs.slide_width, prs.slide_height = W, H

# 01 — Ouverture
s = base_slide(prs, 1, "Ouverture", accent=GREEN)
add_glow(s, Inches(8.8), Inches(-1.5), Inches(6.2), GREEN, 88)
add_glow(s, Inches(9.9), Inches(2.6), Inches(4.6), BLUE, 90)
add_rect(s, Inches(0), Inches(0), Inches(0.16), H, fill=GREEN, radius=False)
add_text(s, Inches(0.92), Inches(0.78), Inches(5.2), Inches(0.32),
         "THÈSE PROFESSIONNELLE · DATA & IA", size=13, color=GREEN_2, bold=True)
add_text(s, Inches(0.92), Inches(1.45), Inches(10.7), Inches(2.05),
         "Rendre un procédé d’extrusion\nlisible, comparable et prédictible",
         size=38, color=WHITE, bold=True, font=FONT_DISPLAY, lines=0.95)
add_text(s, Inches(0.94), Inches(3.78), Inches(8.8), Inches(0.9),
         "Jumeau numérique et IA explicable pour l’extrusion bivis\n"
         "de composants de batteries tout-solide",
         size=19, color=MUTED, lines=1.2)
add_rect(s, Inches(9.65), Inches(4.36), Inches(2.45), Inches(0.86),
         fill=WHITE, line=None)
s.shapes.add_picture(str(ASSETS / "rondol_logo.png"), Inches(9.84), Inches(4.53),
                     width=Inches(2.08))
add_text(s, Inches(0.94), Inches(5.72), Inches(8.4), Inches(0.86),
         "Wilfried Galtier MBEUMI  ·  Nexa Digital School\n"
         "Encadrement industriel : Maël Gallas  ·  2025–2026",
         size=15, color=WHITE, lines=1.35)
note(s, """Bonjour. Je m'appelle Wilfried Galtier MBEUMI. Pendant cette soutenance,
je vais montrer comment j'ai transformé un problème industriel très concret en un
outil démontrable : rendre un procédé d'extrusion bivis lisible, comparable et
prédictible. Ces trois verbes structurent le projet. Lisible, parce que l'opérateur
doit comprendre l'état du procédé. Comparable, parce qu'une configuration de vis
doit pouvoir être rapprochée d'une autre. Prédictible, parce qu'il faut détecter une
dérive avant qu'elle ne devienne un défaut matière. Je présenterai aussi un résultat
qui semblait excellent, puis l'audit qui a montré qu'il était artificiellement
optimiste. Cette correction est devenue l'apport méthodologique central du mémoire.
[Durée cible : 1 minute]""")

# 02 — Verrou industriel
s = base_slide(prs, 2, "Contexte", accent=AMBER)
title(s, "La chimie progresse. La mise en forme reste le verrou.",
      "Batteries tout-solide · extrusion sèche / semi-sèche", accent=AMBER)
card(s, M, Inches(2.05), Inches(3.78), Inches(3.1),
     "01  Voie humide", "Solvant NMP toxique, séchage énergivore et chaîne de production longue.",
     accent=RED)
card(s, Inches(4.78), Inches(2.05), Inches(3.78), Inches(3.1),
     "02  Extrusion sèche", "Procédé continu, moins de solvants et savoir-faire historique de Rondol.",
     accent=GREEN)
card(s, Inches(8.84), Inches(2.05), Inches(3.78), Inches(3.1),
     "03  Verrou actuel", "Formulations céramiques abrasives, peu de données et stabilité difficile à prévoir.",
     accent=AMBER)
add_rich(s, M, Inches(5.55), Inches(11.9), Inches(0.8), [
    ("Marché aval : ", MUTED, False), ("0,26 → 1,77 Md USD", WHITE, True),
    (" d’ici 2031 · TCAC ", MUTED, False), ("37,5 %", AMBER, True),
    (" — un secteur encore pré-industriel.", MUTED, False)], size=18,
    align=PP_ALIGN.CENTER)
note(s, """Le contexte tient en trois idées. Premièrement, les batteries tout-solide
promettent davantage de densité énergétique et une meilleure sécurité, mais leur
industrialisation bute sur la fabrication des électrodes et des électrolytes.
Deuxièmement, la voie humide utilise encore la NMP et impose une étape de séchage
coûteuse. Troisièmement, l'extrusion bivis sèche ou semi-sèche constitue une voie
continue et plus sobre, qui correspond au métier de Rondol. Mais les formulations
SSB sont très chargées, abrasives et encore peu documentées. Les estimations de
marché divergent fortement : cette divergence ne doit pas être masquée, elle traduit
la volatilité du calendrier d'industrialisation. Le sujet de la thèse est donc un
sujet de procédé et d'aide à la décision, pas simplement un sujet d'algorithme.
[Durée cible : 1 minute 30]""")

# 03 — Problème concret
s = base_slide(prs, 3, "Contexte", accent=RED)
title(s, "Chez Rondol, l’opérateur arbitre dans un espace immense.",
      "Extrudeuse bivis Ø 10,5 mm · campagne d’essais d’avril 2026", accent=RED)
xs = [M, Inches(3.72), Inches(6.72), Inches(9.72)]
vals = [("81", "positions de vis", "à configurer"),
        ("13", "types d’éléments", "aux effets différents"),
        ("12", "capteurs", "utilisés pour surveiller"),
        ("8", "essais exploitables", "pour entraîner et tester")]
for x, (v, lab, sub) in zip(xs, vals):
    metric(s, x, Inches(2.35), Inches(2.55), v, lab, accent=CYAN, sub=sub)
add_rect(s, M, Inches(4.55), Inches(11.9), Inches(1.38), fill=RGBColor(0x28, 0x19, 0x1E), line=RED)
add_text(s, Inches(1.0), Inches(4.70), Inches(11.3), Inches(1.08),
         "L’expérience permet de régler. Elle ne permet ni de comparer systématiquement,\n"
         "ni d’anticiper une dérive.", size=21, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, lines=1.1)
note(s, """Voici le problème tel que je l'ai rencontré. Une vis compte 81 positions et
la bibliothèque machine contient 13 types d'éléments. L'espace combinatoire est trop
grand pour être exploré mentalement. Douze capteurs de température enregistrent le
procédé, mais ils servent principalement à constater l'état courant. Enfin, seulement
huit essais étaient exploitables pour la modélisation. Le problème n'est pas que
l'opérateur manque de savoir-faire : au contraire, son expérience est indispensable.
Le problème est que cette expérience reste difficile à formaliser, à comparer et à
transmettre. L'outil devait donc renforcer la décision humaine, sans prétendre la
remplacer. Le faible nombre d'essais deviendra ensuite la limite structurante du volet
prédictif.
[Durée cible : 1 minute 30]""")

# 04 — Problématique
s = base_slide(prs, 4, "Cadrage", accent=CYAN)
title(s, "Une question, trois exigences non négociables.", accent=CYAN)
add_rect(s, M, Inches(1.95), Inches(11.9), Inches(1.48), fill=PANEL_2, line=CYAN)
add_text(s, Inches(1.03), Inches(2.08), Inches(11.25), Inches(1.20),
         "Comment rendre un procédé d’extrusion lisible, comparable et prédictible\n"
         "— sans jamais devenir une boîte noire ?", size=23, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
w3, g3 = three_columns()
for i, (h, b, a) in enumerate([
    ("TRAÇABLE", "Toute valeur doit être expliquée à un ingénieur procédé.", BLUE),
    ("HONNÊTE", "Ce qui n’est pas calibré industriellement est annoncé.", AMBER),
    ("DÉMONTRABLE", "Un outil utilisable devant le client, pas une maquette.", GREEN),
]):
    x = M + i * (w3 + g3)
    card(s, x, Inches(4.05), w3, Inches(2.05), h, b, accent=a, body_size=17)
note(s, """La problématique combine trois ambitions, mais elle impose surtout trois
garde-fous. Traçable : une valeur affichée doit pouvoir être reliée à sa source et à
son calcul. Honnête : une équation non calibrée n'est pas présentée comme une vérité
industrielle. Démontrable : le résultat doit fonctionner devant le client, avec les
contraintes réelles de navigation, d'état et de persistance. Le terme boîte noire est
important. Le modèle statistique n'a pas le droit de produire seul une décision
opérationnelle. Il fournit un indicateur. L'agent métier formule une recommandation
explicable. L'humain conserve l'arbitrage. Ces exigences ont guidé l'architecture et
servent aussi de critères d'évaluation de la solution.
[Durée cible : 1 minute 15]""")

# 05 — Architecture
s = base_slide(prs, 5, "Solution", accent=GREEN)
title(s, "Une architecture en couches — une seule source de vérité.",
      "Le calcul procédé est exécuté une fois, puis réutilisé", accent=GREEN)
layers = [
    ("HMI INDUSTRIELLE", "7 pages Streamlit · opérateur", CYAN),
    ("PERSISTANCE", "édition → appliqué → historique", GREEN),
    ("AGENT EXPLICABLE", "alertes · preuves · recommandations", AMBER),
    ("MOTEUR PROCÉDÉ", "nœuds → zones → machine", BLUE),
    ("BACKBONE GÉOMÉTRIQUE", "Network 7 · 81 positions", GREEN),
]
for i, (h, b, a) in enumerate(layers):
    y = Inches(1.98 + i * 0.88)
    w = Inches(8.25 - i * 0.30)
    x = M + Inches(i * 0.15)
    add_rect(s, x, y, w, Inches(0.72), fill=PANEL if i < 4 else RGBColor(0x0D, 0x32, 0x25), line=a)
    add_text(s, x + Inches(0.26), y + Inches(0.13), Inches(3.25), Inches(0.3), h,
             size=16, color=a, bold=True)
    add_text(s, x + Inches(3.5), y + Inches(0.13), w - Inches(3.75), Inches(0.3), b,
             size=16, color=WHITE, align=PP_ALIGN.RIGHT)
add_rect(s, Inches(9.45), Inches(2.18), Inches(3.17), Inches(3.55),
         fill=RGBColor(0x0C, 0x36, 0x26), line=GREEN)
add_text(s, Inches(9.72), Inches(2.66), Inches(2.63), Inches(0.45),
         "PRINCIPE FONDATEUR", size=14, color=GREEN_2, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(9.72), Inches(3.35), Inches(2.63), Inches(1.42),
         "ENVELOPPER\n≠\nRECALCULER", size=23, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, lines=1.15)
add_text(s, Inches(9.72), Inches(5.04), Inches(2.63), Inches(0.38),
         "1 appel à Network 7", size=14, color=MUTED, align=PP_ALIGN.CENTER)
note(s, """L'architecture répond directement au risque de contradiction entre les
pages. Le backbone géométrique Network 7 est la source de vérité pour les volumes,
le remplissage et le temps de séjour. Le moteur procédé l'appelle une seule fois et
enveloppe son résultat dans des états par position, par zone puis par machine. L'agent
et l'interface consomment ces états ; ils ne recalculent pas la physique chacun de
leur côté. Au-dessus, la persistance sépare trois niveaux : ce que l'opérateur édite,
ce qu'il a réellement validé et l'historique. Ce principe évite qu'une modification
en cours change silencieusement la supervision. La cohérence n'est donc pas obtenue
par une vérification a posteriori : elle est imposée par l'architecture.
[Durée cible : 1 minute 45]""")

# 06 — Cœur métier
s = base_slide(prs, 6, "Solution", accent=BLUE)
title(s, "La vis devient un objet calculable et comparable.",
      "Géométrie réelle Rondol · 81 positions", accent=BLUE)
picture_contain(s, FIG / "fig_engine_logic.png", M, Inches(1.95), Inches(8.15), Inches(4.42))
card(s, Inches(9.18), Inches(1.95), Inches(3.44), Inches(1.22),
     "POSITION", "Élément, zone, matériau et état local.", accent=BLUE, body_size=15)
card(s, Inches(9.18), Inches(3.33), Inches(3.44), Inches(1.22),
     "ZONE", "Remplissage, résidence et couple agrégés.", accent=CYAN, body_size=15)
card(s, Inches(9.18), Inches(4.71), Inches(3.44), Inches(1.22),
     "MACHINE", "KPIs cohérents et historisables.", accent=GREEN, body_size=15)
note(s, """Le cœur métier est la vis, position par position. Chaque position connaît
son élément, sa zone thermique, le matériau nominal et les grandeurs issues du calcul
procédé. Ces positions sont ensuite agrégées sans réinventer les totaux. Cela permet
de comparer deux configurations de manière explicite : où se situe le convoyage, où
se concentre le malaxage, où apparaît une restriction et comment cela affecte le
remplissage ou le temps de séjour. Cette granularité est indispensable pour produire
des recommandations localisées. Une alerte qui dit seulement que le procédé est
instable est peu utile. Une alerte qui identifie la zone et la configuration qui y
contribuent devient actionnable pour l'opérateur.
[Durée cible : 1 minute 30]""")

# 07 — Physique
s = base_slide(prs, 7, "Solution", accent=AMBER)
title(s, "La physique est embarquée — et ses limites restent visibles.",
      "Modèle nominal, documenté, non présenté comme calibration industrielle", accent=AMBER)
w3, g3 = three_columns()
physics = [
    ("VISCOSITÉ", "η(γ̇, T)", "Carreau–Yasuda + Arrhenius", BLUE),
    ("COUPLE LOCAL", "M = η·γ̇²·V / 2πN", "Calcul nœud par nœud", GREEN),
    ("THERMIQUE", "Tᵣ = Tₛ + P/ṁCₚ + kτ", "Équation encadrement industriel", AMBER),
]
for i, (h, formula, b, a) in enumerate(physics):
    x = M + i * (w3 + g3)
    add_rect(s, x, Inches(2.12), w3, Inches(2.68), fill=PANEL, line=a)
    add_text(s, x + Inches(0.25), Inches(2.42), w3 - Inches(0.5), Inches(0.32), h,
             size=15, color=a, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(3.02), w3 - Inches(0.4), Inches(0.7), formula,
             size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.25), Inches(4.02), w3 - Inches(0.5), Inches(0.36), b,
             size=15, color=MUTED, align=PP_ALIGN.CENTER)
add_rect(s, M, Inches(5.25), Inches(11.9), Inches(0.9), fill=RGBColor(0x2A, 0x22, 0x13), line=AMBER)
add_text(s, Inches(0.98), Inches(5.45), Inches(11.35), Inches(0.45),
         "E5 / E6 / E7 restent différées : l’interface affiche « À venir » plutôt qu’une valeur plausible.",
         size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
note(s, """Trois familles de calculs sont effectivement embarquées : la viscosité
locale, le couple local et une relation thermique nominale fournie par l'encadrement.
Elles servent à structurer la lecture du procédé et non à annoncer une précision que
nous ne possédons pas encore. Les équations E5, E6 et E7, notamment l'énergie locale,
la température avancée et la pression filière, restent volontairement différées.
Dans l'application, elles renvoient None et sont affichées comme à venir. C'est un
choix important : remplir une carte avec une valeur plausible aurait amélioré la
démonstration visuelle, mais dégradé la crédibilité scientifique. Le prochain jalon
est la calibration sur des campagnes instrumentées, pas l'ajout d'équations non
validées.
[Durée cible : 1 minute 30]""")

# 08 — Agent explicable
s = base_slide(prs, 8, "Solution", accent=GREEN)
title(s, "Le modèle prédit. Les règles expliquent. L’humain décide.",
      "Deux niveaux d’intelligence, volontairement séparés", accent=GREEN)
steps = [
    ("ÉTAT PROCÉDÉ", "capteurs + moteur", BLUE),
    ("RÈGLES MÉTIER", "seuils explicites", AMBER),
    ("ALERTES", "preuve chiffrée", RED),
    ("ACTION", "recommandation localisée", GREEN),
    ("OPÉRATEUR", "arbitrage final", CYAN),
]
sw = Inches(2.1)
for i, (h, b, a) in enumerate(steps):
    x = Inches(0.55 + i * 2.55)
    add_rect(s, x, Inches(2.45), sw, Inches(2.15), fill=PANEL, line=a)
    add_text(s, x + Inches(0.16), Inches(2.83), sw - Inches(0.32), Inches(0.38), h,
             size=15, color=a, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.16), Inches(3.42), sw - Inches(0.32), Inches(0.68), b,
             size=16, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        arrow(s, x + sw, Inches(3.52), x + Inches(2.5), Inches(3.52), color=MUTED)
add_text(s, M, Inches(5.4), Inches(11.9), Inches(0.68),
         "Une recommandation cite toujours sa preuve, sa zone et le seuil franchi.",
         size=22, color=GREEN_2, bold=True, align=PP_ALIGN.CENTER)
note(s, """L'agent n'est pas un classifieur déguisé. Il applique des règles métier
explicites à l'état du procédé. Chaque alerte cite la variable, la valeur, le seuil
et la zone concernée. La recommandation reste donc vérifiable. Le modèle statistique
intervient à un autre niveau : il produit un score de stabilité et une probabilité
de dérive. Il ne commande jamais directement une action. Cette séparation répond à
la problématique sans boîte noire et réduit le risque d'automatisation aveugle. Dans
la démonstration, je peux montrer comment une modification de configuration change
les alertes métier immédiatement, tandis que le score ML reste clairement présenté
comme un indicateur expérimental issu de la campagne d'avril.
[Durée cible : 1 minute 30]""")

# 09 — Données
s = base_slide(prs, 9, "Données", accent=CYAN)
title(s, "Une campagne réelle — et une faiblesse impossible à masquer.",
      "7–13 avril 2026 · 12 fichiers CSV · capteurs de température", accent=CYAN)
vals = [("310 782", "relevés bruts", "25,8 Mo"),
        ("10–16 %", "couverture par capteur", "acquisition fragmentée"),
        ("627", "fenêtres de 60 s", "après nettoyage"),
        ("87", "variables du modèle", "12 × 7 + 3 croisées")]
for i, (v, l, sub) in enumerate(vals):
    metric(s, M + i * Inches(3.0), Inches(2.15), Inches(2.72), v, l, accent=CYAN, sub=sub)
add_rect(s, M, Inches(4.45), Inches(7.68), Inches(1.45), fill=PANEL, line=GRID)
add_text(s, Inches(0.98), Inches(4.75), Inches(7.15), Inches(0.75),
         "Nettoyage : codes thermocouple, doublons, resampling 10 s,\nforward-fill borné à 60 s.",
         size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(s, Inches(8.68), Inches(4.45), Inches(3.94), Inches(1.45), fill=RGBColor(0x2A, 0x22, 0x13), line=AMBER)
add_text(s, Inches(8.95), Inches(4.68), Inches(3.4), Inches(0.95),
         "8 essais seulement\n→ limite structurelle", size=22, color=AMBER,
         bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
note(s, """Les données sont réelles. Douze capteurs produisent 310 782 relevés bruts,
mais l'acquisition est fragmentée : chaque capteur ne couvre qu'une partie de la
campagne. Le nettoyage a supprimé les codes thermocouple aberrants, traité les
doublons et remis les séries sur une grille de dix secondes avec un remplissage
strictement borné. Après fenêtrage, 627 fenêtres réelles de soixante secondes sont
exploitables. Les 87 variables correspondent exactement aux entrées attendues par le
modèle déployé : sept statistiques sur douze voies, plus trois gradients croisés.
La limite majeure reste le nombre d'essais : huit. Toute la stratégie d'évaluation
doit donc empêcher que les nombreuses fenêtres d'un même essai donnent l'illusion
d'une grande diversité expérimentale.
[Durée cible : 1 minute 45]""")

# 10 — Validation
s = base_slide(prs, 10, "Modélisation", accent=BLUE)
title(s, "La séparation par essai change la vérité du modèle.",
      "Leave-One-Group-Out · un essai entier exclu à chaque pli", accent=BLUE)
metric(s, Inches(1.0), Inches(2.15), Inches(4.75), "0,92", "validation naïve",
       accent=RED, sub="fenêtres mélangées · optimiste")
metric(s, Inches(7.58), Inches(2.15), Inches(4.75), "0,79", "validation stricte",
       accent=GREEN, sub="essais séparés · valeur réaliste")
arrow(s, Inches(5.95), Inches(2.95), Inches(7.35), Inches(2.95), color=AMBER, width=3)
add_text(s, Inches(5.8), Inches(3.38), Inches(1.8), Inches(0.42), "− 15 pts",
         size=19, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
add_rect(s, M, Inches(4.45), Inches(11.9), Inches(1.35), fill=PANEL, line=GRID)
add_rich(s, Inches(1.0), Inches(4.72), Inches(11.35), Inches(0.82), [
    ("87,5 %", CYAN, True), (" des fenêtres utilisées pour l’entraînement en moyenne · ", MUTED, False),
    ("aucun essai", WHITE, True), (" simultanément dans train et test.", MUTED, False)],
    size=20, align=PP_ALIGN.CENTER)
note(s, """Le protocole de validation est plus important que le choix de l'algorithme.
Les fenêtres d'un même essai sont autocorrélées. Si elles sont réparties au hasard,
le modèle retrouve au test des motifs presque identiques à ceux vus à l'entraînement.
Le score monte alors à environ 0,92. Avec une séparation stricte par essai, il tombe
à environ 0,79. Cet écart de quinze points n'est pas une dégradation provoquée par le
protocole ; c'est la mesure de l'optimisme que le protocole naïf cachait. Le
Leave-One-Group-Out exclut chaque essai à tour de rôle et entraîne sur 87,5 % des
fenêtres en moyenne. Sept plis sur huit dépassent le seuil de 70 % demandé par le
référentiel ; l'exception vient du poids disproportionné de l'essai 1.
[Durée cible : 2 minutes]""")

# 11 — Résultat initial
s = base_slide(prs, 11, "Audit scientifique", accent=GREEN)
title(s, "Le résultat semblait décisif.",
      "Augmentation globale · 800 fenêtres synthétiques", accent=GREEN)
add_text(s, Inches(1.15), Inches(2.0), Inches(4.0), Inches(0.4), "SANS AUGMENTATION",
         size=14, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.15), Inches(2.0), Inches(4.0), Inches(0.4), "AVEC AUGMENTATION",
         size=14, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.15), Inches(2.65), Inches(4.0), Inches(1.0), "0,809",
         size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_DISPLAY)
add_text(s, Inches(8.15), Inches(2.65), Inches(4.0), Inches(1.0), "0,918",
         size=56, color=GREEN_2, bold=True, align=PP_ALIGN.CENTER, font=FONT_DISPLAY)
arrow(s, Inches(5.2), Inches(3.14), Inches(7.9), Inches(3.14), color=GREEN, width=4)
chip(s, Inches(5.72), Inches(3.52), Inches(1.65), "+0,109", fill=RGBColor(0x0C, 0x36, 0x26), color=GREEN_2)
add_text(s, Inches(1.2), Inches(4.35), Inches(10.9), Inches(0.62),
         "Variance inter-essais divisée par plus de trois : 0,176 → 0,054",
         size=19, color=MUTED, align=PP_ALIGN.CENTER)
add_rect(s, M, Inches(5.38), Inches(11.9), Inches(0.92), fill=RGBColor(0x2D, 0x18, 0x1F), line=RED)
add_text(s, Inches(0.98), Inches(5.64), Inches(11.35), Inches(0.38),
         "J’ai présenté ce résultat comme l’aboutissement du volet prédictif. Il était faux.",
         size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
note(s, """Pour compenser le faible nombre d'essais, j'ai généré un pool synthétique
par bootstrap conditionné par classe, avec un bruit borné et la reproduction de
certaines imperfections capteur. Le résultat semblait excellent : le Random Forest
passait de 0,809 à 0,918 et sa variance inter-essais était divisée par plus de trois.
À ce stade, tout semblait cohérent : les données synthétiques n'étaient ajoutées
qu'au train, jamais au test. J'ai donc présenté ce gain comme l'aboutissement du
volet prédictif. Il faut marquer une pause sur cette diapositive, parce que la suite
change l'interprétation du projet : le score était réel dans le fichier de résultats,
mais le protocole qui l'avait produit était contaminé.
[Durée cible : 1 minute 30]""")

# 12 — Fuite par ancrage
s = base_slide(prs, 12, "Audit scientifique", accent=RED, dark=RGBColor(0x16, 0x0D, 0x14))
title(s, "La fuite était en amont du pli de test.",
      "Fuite par ancrage · invisible dans le split final", accent=RED)
nodes = [
    ("8 ESSAIS RÉELS", BLUE),
    ("POOL SYNTHÉTIQUE\nGLOBAL", RED),
    ("TRAIN DU PLI", AMBER),
    ("TEST : 1 ESSAI", GREEN),
]
xs = [Inches(0.65), Inches(3.75), Inches(7.0), Inches(10.15)]
for x, (lab, a) in zip(xs, nodes):
    add_rect(s, x, Inches(2.5), Inches(2.5), Inches(1.35), fill=PANEL, line=a)
    add_text(s, x + Inches(0.15), Inches(2.82), Inches(2.2), Inches(0.65), lab,
             size=17, color=a, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)
for i in range(3):
    arrow(s, xs[i] + Inches(2.5), Inches(3.18), xs[i + 1] - Inches(0.12), Inches(3.18),
          color=MUTED)
# Retour contaminant du test vers le pool global.
c = s.shapes.add_connector(MSO_CONNECTOR.ELBOW, Inches(11.4), Inches(4.0), Inches(5.0), Inches(4.85))
c.line.color.rgb = RED
c.line.width = Pt(3)
c.line.dash_style = 2
c.line.end_arrowhead = True
add_text(s, Inches(5.0), Inches(4.42), Inches(5.6), Inches(0.72),
         "L’essai exclu avait servi de point d’ancrage au bootstrap.",
         size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)
add_rect(s, M, Inches(5.45), Inches(11.9), Inches(0.82), fill=RGBColor(0x31, 0x18, 0x20), line=RED)
add_text(s, Inches(0.98), Inches(5.67), Inches(11.35), Inches(0.35),
         "Correction : régénérer le pool dans chaque pli, uniquement depuis les essais d’entraînement.",
         size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
note(s, """L'audit a montré une fuite par ancrage. Le pool synthétique était créé une
seule fois à partir des huit essais réels, puis réutilisé dans chaque pli. L'essai
censé être totalement exclu contribuait donc indirectement au train : ses fenêtres
servaient de points d'ancrage au bootstrap et ses valeurs participaient au calcul
des amplitudes de bruit. Le défaut était difficile à voir, parce que le jeu de test
ne contenait aucune ligne synthétique. Le split final semblait propre. La fuite se
trouvait en amont, au moment de la génération. Pour la corriger, j'ai conservé le
même algorithme, le même volume et la même graine, mais j'ai régénéré le pool dans
chaque pli à partir des seuls essais d'entraînement.
[Durée cible : 2 minutes 30]""")

# 13 — Résultat corrigé
s = base_slide(prs, 13, "Audit scientifique", accent=RED)
title(s, "Une fois la fuite retirée, le gain disparaît.",
      "Random Forest · F1-macro Leave-One-Group-Out", accent=RED)
protocols = [
    ("SANS AUGMENTATION", "0,809", "référence", BLUE),
    ("GLOBALE — FUITÉE", "0,918", "+0,109 artificiel", RED),
    ("PAR PLI — CORRIGÉE", "0,809", "−0,001 réel", GREEN),
]
w3, g3 = three_columns()
for i, (h, v, sub, a) in enumerate(protocols):
    x = M + i * (w3 + g3)
    add_rect(s, x, Inches(2.05), w3, Inches(2.82), fill=PANEL, line=a)
    add_text(s, x + Inches(0.2), Inches(2.42), w3 - Inches(0.4), Inches(0.35), h,
             size=14, color=a, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(3.0), w3 - Inches(0.3), Inches(0.88), v,
             size=48, color=WHITE if i != 1 else RED, bold=True,
             align=PP_ALIGN.CENTER, font=FONT_DISPLAY)
    add_text(s, x + Inches(0.2), Inches(4.08), w3 - Inches(0.4), Inches(0.38), sub,
             size=17, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
add_rich(s, M, Inches(5.35), Inches(11.9), Inches(0.75), [
    ("Les cinq modèles restent entre ", MUTED, False), ("0,781 et 0,824", WHITE, True),
    (" : aucun ne se détache statistiquement.", MUTED, False)], size=20,
    align=PP_ALIGN.CENTER)
note(s, """Après correction, le Random Forest revient à 0,809. Le gain passe de plus
0,109 à moins 0,001 : il disparaît. La conclusion ne concerne pas seulement le
Random Forest. Les cinq familles de modèles testées restent groupées entre 0,781 et
0,824, avec des écarts-types importants. Il serait donc abusif de déclarer un
vainqueur statistique. Le SVM obtient le score moyen le plus élevé, mais l'écart est
trop faible au regard de la variabilité. Le Random Forest reste déployé pour des
raisons opérationnelles : interprétabilité, tolérance aux valeurs manquantes et
stabilité inter-essais. Le choix du modèle devient un compromis d'ingénierie, pas un
classement fondé sur la troisième décimale.
[Durée cible : 2 minutes]""")

# 14 — Valeur de l'audit
s = base_slide(prs, 14, "Audit scientifique", accent=CYAN)
title(s, "Dire la vérité renforce le projet.",
      "Le résultat corrigé vaut davantage qu’un score flatteur", accent=CYAN)
w3, g3 = three_columns()
items = [
    ("DÉONTOLOGIQUE", "Ne pas promettre à Rondol une performance qui ne tiendrait pas en production.", RED),
    ("MÉTHODOLOGIQUE", "Montrer que la performance dépend autant du protocole que de l’algorithme.", CYAN),
    ("REPRODUCTIBLE", "Livrer prédictions, métriques et générateur afin que chaque chiffre soit contestable.", GREEN),
]
for i, (h, b, a) in enumerate(items):
    card(s, M + i * (w3 + g3), Inches(2.18), w3, Inches(3.35), h, b,
         accent=a, body_size=18)
add_text(s, M, Inches(5.92), Inches(11.9), Inches(0.42),
         "Ce n’est pas l’échec du modèle : c’est la réussite de l’audit.",
         size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
note(s, """J'expose cette correction pour trois raisons. D'abord, une raison
déontologique : annoncer 0,918 aurait créé une attente que le modèle ne pouvait pas
tenir. Ensuite, une raison méthodologique : cet exemple montre concrètement qu'un
score n'a de sens qu'avec son protocole. Enfin, une raison de reproductibilité : les
prédictions par pli, les métriques et le générateur corrigé sont livrés. Un tiers peut
donc recalculer et contester chaque valeur. Cet épisode renforce le projet parce qu'il
transforme une faiblesse cachée en connaissance explicite. Il donne aussi une règle
pour la suite : aucune augmentation ne doit être générée avant la séparation des
groupes, et les contrôles méthodologiques doivent devenir des tests permanents.
[Durée cible : 1 minute 30]""")

# 15 — Validation externe
s = base_slide(prs, 15, "Validation", accent=AMBER)
title(s, "Le modèle conserve un signal — pas une certification.",
      "Validation externe sans réentraînement", accent=AMBER)
metric(s, M, Inches(2.15), Inches(3.55), "3 479", "fenêtres évaluées",
       accent=CYAN, sub="15 runs simulés")
metric(s, Inches(4.88), Inches(2.15), Inches(3.55), "0,753", "AUC externe",
       accent=AMBER, sub="pouvoir discriminant conservé")
metric(s, Inches(9.04), Inches(2.15), Inches(3.55), "62 %", "instabilités détectées",
       accent=GREEN, sub="erreurs plutôt conservatrices")
add_rect(s, M, Inches(4.55), Inches(11.9), Inches(1.18), fill=PANEL, line=GRID)
add_text(s, Inches(1.0), Inches(4.72), Inches(11.3), Inches(0.86),
         "L’écart mesure la sensibilité au changement de distribution.\n"
         "Le modèle reste un indicateur d’aide à la décision — pas un détecteur certifié.",
         size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER, lines=1.15)
note(s, """Pour éprouver la transférabilité sans attendre une nouvelle campagne, j'ai
appliqué le modèle sans réentraînement à une base continue générée à partir de
l'échantillon. L'AUC est de 0,753 sur 3 479 fenêtres et quinze runs. Le signal
discriminant subsiste, mais la baisse confirme une sensibilité au changement de
distribution. Environ 62 % des fenêtres instables sont détectées. Les erreurs sont
majoritairement conservatrices : davantage de fausses alertes que de dérives manquées.
Je ne présente pas cette base comme une preuve industrielle. Elle sert à tester le
comportement du pipeline hors du jeu d'entraînement et confirme le bon niveau de
promesse : un indicateur expérimental d'aide à la décision, à recalibrer sur de
nouvelles campagnes réelles.
[Durée cible : 1 minute 30]""")

# 16 — Application
s = base_slide(prs, 16, "Produit", accent=GREEN)
title(s, "Une HMI industrielle — pas un tableau de bord générique.",
      "Supervision · profil de vis · analyse · historique · moteur procédé", accent=GREEN)
picture_contain(s, CAP / "cap_supervision.png", M, Inches(1.86), Inches(9.15), Inches(4.9),
                background=DARK, pad=0.05, line=GREEN)
card(s, Inches(10.16), Inches(1.86), Inches(2.46), Inches(1.18), "7 PAGES",
     "Parcours opérateur cohérent.", accent=CYAN, body_size=14)
card(s, Inches(10.16), Inches(3.18), Inches(2.46), Inches(1.18), "3 ÉTATS",
     "Édition · appliqué · historique.", accent=GREEN, body_size=14)
card(s, Inches(10.16), Inches(4.50), Inches(2.46), Inches(1.18), "FR / EN",
     "Interface bilingue et testée.", accent=AMBER, body_size=14)
chip(s, Inches(10.16), Inches(6.02), Inches(2.46), "WCAG 2.1 AA", fill=RGBColor(0x0C, 0x36, 0x26), color=GREEN_2)
note(s, """La solution se matérialise dans une HMI sombre inspirée des interfaces
industrielles de Rondol. La page Supervision rassemble le score expérimental, la
probabilité de dérive, les alertes explicables, les recommandations et les KPIs
procédé. Les sept pages couvrent la configuration de vis, les réglages, l'analyse de
run, l'historique, le moteur procédé et le compte utilisateur. La persistance en
trois couches évite qu'une édition non enregistrée modifie la supervision. L'interface
est bilingue et les contrastes ont été mesurés entre 6,38:1 et 18,39:1. Quinze tests
spécifiques verrouillent les alternatives textuelles, les éléments décoratifs et les
contrastes. Une démonstration en direct peut compléter cette diapositive.
[Durée cible : 2 minutes, démonstration non comprise]""")

# 17 — Cas de démonstration
s = base_slide(prs, 17, "Produit", accent=GREEN)
title(s, "Cinq cas prouvent la sensibilité et la réversibilité.",
      "C1 → C5 · défaut provoqué, recommandation, correction", accent=GREEN)
cases = [
    ("C1", "65", "référence", BLUE),
    ("C2", "82", "optimisée", GREEN),
    ("C3", "46", "défaut Z5", RED),
    ("C4", "→", "recommandation", AMBER),
    ("C5", "+32", "alerte levée", GREEN),
]
cw, gap = Inches(2.13), Inches(0.25)
for i, (code, val, lab, a) in enumerate(cases):
    x = M + i * (cw + gap)
    add_rect(s, x, Inches(2.18), cw, Inches(3.05), fill=PANEL, line=a)
    add_text(s, x, Inches(2.5), cw, Inches(0.35), code, size=17, color=a,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.14), cw, Inches(0.8), val, size=38, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER, font=FONT_DISPLAY)
    add_text(s, x + Inches(0.12), Inches(4.23), cw - Inches(0.24), Inches(0.45), lab,
             size=15, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
add_rich(s, M, Inches(5.68), Inches(11.9), Inches(0.6), [
    ("C3 → C5 : ", GREEN_2, True), ("+32 points · +0,52 de probabilité · alerte levée", WHITE, True),
    (" — la projection de C4 se vérifie.", MUTED, False)], size=19,
    align=PP_ALIGN.CENTER)
note(s, """Les cinq cas démontrent trois propriétés. La sensibilité : entre C1 et C2,
une configuration améliorée fait passer le score de 65 à 82. La détectabilité : en
C3, un défaut provoqué fait chuter le score à 46 et localise une alerte rouge en
zone 5. L'actionnabilité : en C4, l'agent propose une correction chiffrée. Enfin, la
réversibilité : après application en C5, le score remonte de 32 points, la probabilité
de stabilité augmente de 0,52 et l'alerte disparaît. Le point le plus intéressant
est que la projection annoncée par l'agent en C4 se vérifie dans C5. Le scénario ne
montre donc pas seulement une interface réactive ; il montre une chaîne complète de
diagnostic, recommandation et contrôle après correction.
[Durée cible : 1 minute 45]""")

# 18 — Pilotage
s = base_slide(prs, 18, "Pilotage", accent=BLUE)
title(s, "Le projet a été piloté par validation, pas par accumulation.",
      "Kanban · WIP = 1 · jalons et risques tracés", accent=BLUE)
picture_contain(s, FIG / "fig_gantt.png", M, Inches(1.92), Inches(8.55), Inches(4.52))
card(s, Inches(9.55), Inches(1.92), Inches(3.07), Inches(1.25), "WIP = 1",
     "Une seule brique en cours avant validation.", accent=CYAN, body_size=14)
card(s, Inches(9.55), Inches(3.34), Inches(3.07), Inches(1.25), "2 JALONS",
     "Campagne d’avril et démo client du 16 juin tenues.", accent=GREEN, body_size=14)
card(s, Inches(9.55), Inches(4.76), Inches(3.07), Inches(1.25), "12 INCIDENTS",
     "Cause, correctif et test de non-régression.", accent=AMBER, body_size=14)
note(s, """Le projet a suivi un Kanban à encours limité. Une brique n'était pas
entamée tant que la précédente n'avait pas été validée par l'encadrement. Ce choix a
ralenti certains démarrages mais réduit les reprises et les incohérences. Deux jalons
datés ont structuré le calendrier : la campagne d'essais d'avril et la démonstration
client du 16 juin. Les risques ont été cartographiés, notamment la faible volumétrie,
la perte d'état en environnement cloud, les dépendances entre modules et la tentation
de sur-promettre la physique nominale. Douze incidents techniques ont été tracés avec
leur cause, leur correction et leur référence de commit, puis transformés en tests
de non-régression. Le pilotage fait donc partie du résultat, pas seulement de son
contexte.
[Durée cible : 1 minute 30]""")

# 19 — Qualité
s = base_slide(prs, 19, "Qualité", accent=GREEN)
title(s, "La qualité logicielle protège la crédibilité scientifique.",
      "Tests · accessibilité · performance · persistance", accent=GREEN)
vals = [
    ("720", "tests automatisés", "tous passants", GREEN),
    ("15", "tests accessibilité", "WCAG + ARIA", CYAN),
    ("×1030", "accélération SQL", "index B-tree", AMBER),
    ("12", "incidents figés", "non-régression", RED),
]
for i, (v, l, sub, a) in enumerate(vals):
    metric(s, M + i * Inches(3.0), Inches(2.18), Inches(2.72), v, l, accent=a, sub=sub)
add_rect(s, M, Inches(4.55), Inches(11.9), Inches(1.25), fill=PANEL, line=GREEN)
add_text(s, Inches(1.02), Inches(4.72), Inches(11.25), Inches(0.88),
         "Un résultat scientifique n’est crédible que si le logiciel qui le produit\n"
         "est reproductible, testable et observable.", size=21, color=WHITE,
         bold=True, align=PP_ALIGN.CENTER, lines=1.1)
note(s, """La suite compte 720 tests automatisés, tous passants. Elle couvre la logique
métier, les pages Streamlit, la persistance, l'internationalisation, l'accessibilité
et les incidents de production. Quinze tests concernent directement l'accessibilité.
Le benchmark SQL montre une accélération d'environ 1 030 fois entre une recherche
non indexée et une table équipée d'un index B-tree ; ce chiffre illustre la démarche
de mesure, pas une promesse de capacité industrielle. Douze incidents réels ont été
figés en non-régression. Cette qualité n'est pas décorative : une fuite d'état, une
valeur par défaut silencieuse ou une incohérence entre pages pourrait invalider la
démonstration scientifique. Tester le logiciel revient donc à protéger la validité
du raisonnement.
[Durée cible : 1 minute 30]""")

# 20 — Conclusion
s = base_slide(prs, 20, "Conclusion", accent=GREEN)
add_glow(s, Inches(8.6), Inches(-1.1), Inches(6.5), GREEN, 88)
add_rect(s, Inches(0), Inches(0), Inches(0.16), H, fill=GREEN, radius=False)
add_text(s, Inches(0.95), Inches(0.72), Inches(5.6), Inches(0.35),
         "RÉPONSE À LA PROBLÉMATIQUE", size=13, color=GREEN_2, bold=True)
add_text(s, Inches(0.95), Inches(1.32), Inches(10.7), Inches(1.15),
         "Oui — le procédé devient lisible, comparable\net prédictible. Mais pas autonome.",
         size=34, color=WHITE, bold=True, font=FONT_DISPLAY, lines=1.0)
concl = [
    ("ACQUIS", "Jumeau numérique cohérent · agent explicable · HMI démontrable", GREEN),
    ("ASSUMÉ", "Physique nominale · huit essais · pression filière différée", AMBER),
    ("APPRIS", "La rigueur du protocole vaut davantage qu’un score flatteur", CYAN),
]
for i, (h, b, a) in enumerate(concl):
    y = Inches(3.05 + i * 0.82)
    add_text(s, Inches(0.98), y, Inches(1.5), Inches(0.35), h, size=15, color=a, bold=True)
    add_text(s, Inches(2.35), y, Inches(9.4), Inches(0.42), b, size=18, color=WHITE)
add_rect(s, Inches(0.95), Inches(5.82), Inches(11.45), Inches(0.83), fill=RGBColor(0x0C, 0x36, 0x26), line=GREEN)
add_text(s, Inches(1.18), Inches(6.02), Inches(10.98), Inches(0.42),
         "Le prochain gain viendra de nouvelles campagnes — pas d’un modèle plus sophistiqué.",
         size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
note(s, """La réponse à la problématique est positive, mais encadrée. Le procédé est
plus lisible grâce aux états par position et aux alertes explicables. Il devient
comparable grâce à la géométrie commune et à l'historique. Il devient partiellement
prédictible grâce au score expérimental et au protocole de validation strict. Il ne
devient pas autonome : l'humain reste décisionnaire et la physique n'est pas encore
calibrée industriellement. Ce qui est acquis est démontrable. Ce qui est limité est
écrit. Ce que j'ai surtout appris est qu'un score séduisant ne vaut rien si son
protocole ne résiste pas à l'audit. La prochaine étape n'est donc pas un algorithme
plus sophistiqué, mais de nouvelles campagnes mieux équilibrées, avec davantage de
capteurs et une calibration procédé. Je vous remercie et je suis prêt à répondre à
vos questions.
[Durée cible : 1 minute 30]""")

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"[OK] {OUT}")
print(f"     {len(prs.slides)} diapositives · 16:9")
