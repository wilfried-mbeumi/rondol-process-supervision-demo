# -*- coding: utf-8 -*-
"""Export PPTX -> PDF (PowerPoint COM) + contrôle qualité du support de soutenance."""
from __future__ import annotations
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports"
PPTX = REPORTS / "Soutenance_Rondol_Wilfried_Galtier_MBEUMI_FINAL_PRO.pptx"
PDF = REPORTS / "Soutenance_Rondol_Wilfried_Galtier_MBEUMI_FINAL_PRO.pdf"
SCRIPT_DOCX = REPORTS / "Script_Soutenance_Rondol_Wilfried_Galtier_MBEUMI.docx"
PP_SAVE_PDF = 32


def convert():
    import win32com.client as win32
    ppt = win32.DispatchEx("PowerPoint.Application")
    try:
        ppt.Visible = True
        pres = ppt.Presentations.Open(str(PPTX), WithWindow=False)
        pres.SaveAs(str(PDF), PP_SAVE_PDF)
        pres.Close()
        return True
    except Exception as exc:
        print(f"[ERREUR] PowerPoint COM: {exc}")
        return False
    finally:
        try:
            ppt.Quit()
        except Exception:
            pass


def qc():
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(str(PPTX))
    slides = list(prs.slides)
    n = len(slides)
    n_pics = 0
    n_notes = 0
    full = []
    for sl in slides:
        for sh in sl.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                n_pics += 1
            if sh.has_text_frame:
                full.append(sh.text_frame.text)
        if sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip():
            n_notes += 1
    text = "\n".join(full)
    title = slides[0]
    title_pics = sum(1 for sh in title.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
    checks = {
        "Slides entre 12 et 18": 12 <= n <= 18,
        "Notes orateur sur toutes les slides": n_notes == n,
        "Logos sur la page de titre (>=2 images)": title_pics >= 2,
        "Captures/figures intégrées (>=10 images)": n_pics >= 10,
        "Auteur présent": "Wilfried Galtier MBEUMI" in text,
        "Tuteur Maël Gallas": "Maël Gallas" in text,
        "Référent Moussa NDIAYE": "Moussa NDIAYE" in text,
        "Problématique présente": "Comment concevoir et déployer" in text,
        "PAS de '[INSÉRER'": "[INSÉRER" not in text and "INSÉRER" not in text.upper(),
        "PAS de Markdown brut (##, **)": "##" not in text and "**" not in text,
        "PAS de 'À COMPLÉTER'": "À COMPLÉTER" not in text.upper(),
    }
    print("\n=== CONTRÔLE QUALITÉ SOUTENANCE ===")
    ok = True
    for k, v in checks.items():
        print(f"  [{'OK ' if v else 'KO '}] {k}")
        ok = ok and v
    print(f"  Slides : {n} | Images : {n_pics} | Slides avec notes : {n_notes}")
    return ok, n, n_pics


def pages():
    try:
        from pypdf import PdfReader
        return len(PdfReader(str(PDF)).pages)
    except Exception as exc:
        print(f"[WARN] pages PDF: {exc}")
        return None


def main():
    if not PPTX.exists():
        print("[ERREUR] PPTX introuvable"); return 1
    print("[1/3] Export PDF (PowerPoint COM) ...")
    pdf_ok = convert()
    print(f"      PDF : {'OK ' + str(PDF) if pdf_ok and PDF.exists() else 'ÉCHEC'}")
    print("[2/3] Contrôle qualité ...")
    ok, n, n_pics = qc()
    print("[3/3] Bilan")
    print(f"\nPPTX : {PPTX if PPTX.exists() else 'ABSENT'}")
    print(f"PDF  : {PDF if PDF.exists() else 'ABSENT'}  (pages: {pages()})")
    print(f"Script oral : {SCRIPT_DOCX if SCRIPT_DOCX.exists() else 'ABSENT'}")
    print(f"Slides : {n} | Images : {n_pics}")
    print(f"Qualité globale : {'CONFORME' if ok else 'À VÉRIFIER'}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
