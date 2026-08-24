# -*- coding: utf-8 -*-
"""Vérificateur de synchronisation des livrables de soutenance.

Extrait les chiffres clés de chaque livrable et les compare aux sources de
vérité (JSON de mesures). Sort en erreur (exit 1) au moindre écart.

Usage : python scripts/check_sync_livrables.py
"""
from __future__ import annotations

import json
import re
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
ERREURS: list[str] = []


def err(msg: str) -> None:
    ERREURS.append(msg)


def contient(texte: str, motifs: list[str], source: str, doit: bool = True) -> None:
    for m in motifs:
        ok = m in texte
        if doit and not ok:
            err(f"{source} : motif attendu absent → « {m} »")
        if not doit and ok:
            err(f"{source} : motif interdit présent → « {m} »")


def main() -> int:
    # ---------- sources de vérité ----------
    ml = json.load(open(ROOT / "reports/ml_metrics_w60.json", encoding="utf-8"))
    ext = json.load(open(ROOT / "reports/eval_consolidated_w60.json", encoding="utf-8"))
    gen = json.load(open(ROOT / "data/consolidated/rapport_generation.json", encoding="utf-8"))

    rf = ml["test"]["RandomForest"]
    acc_rf = f"{rf['accuracy']:.3f}".replace(".", ",")          # 0,950
    f1_rf = f"{rf['f1_macro']:.3f}".replace(".", ",")           # 0,917
    auc_ext = f"{ext['roc_auc']:.3f}".replace(".", ",")         # ex. 0,753
    n_rows = f"{gen['n_rows']:,}".replace(",", " ")             # 100 800
    n_win_ext = f"{ext['n_windows']:,}".replace(",", " ")       # 3 479
    TESTS = "725"

    print(f"Vérités : RF acc {acc_rf} / F1 {f1_rf} · ext AUC {auc_ext} ({n_win_ext} fen.) · base {n_rows} · tests {TESTS}")

    # ---------- 1. mémoire (markdown source) ----------
    mem = open(ROOT / "docs/memoire_these_professionnelle_rondol.md", encoding="utf-8").read()
    contient(mem, [f"AUC {auc_ext}", TESTS, "0,918"], "mémoire")
    contient(mem, ["693"], "mémoire", doit=False)
    if re.search(r"\b3 ?479\b", mem) is None:
        err("mémoire : nombre de fenêtres de validation externe absent (3 479)")

    # ---------- 2. README ----------
    rd = open(ROOT / "README.md", encoding="utf-8").read()
    contient(rd, [f"AUC {auc_ext}", "100 800", "798 fenêtres"], "README")

    # ---------- 3. rapport de génération ----------
    rg = open(ROOT / "data/consolidated/rapport_generation.md", encoding="utf-8").read()
    contient(rg, [f"AUC {auc_ext}", "100 800"], "rapport_generation.md")
    pct = f"{gen.get('describe', {}) and ''}"  # % stable vérifié via eval
    contient(rg, [f"{ext['pct_stable']:.1f}".replace(".", ",") + " %"], "rapport_generation.md")

    # ---------- 4. PowerPoint ----------
    try:
        from pptx import Presentation
        prs = Presentation(str(ROOT / "reports/soutenance/MBEUMI_Wilfried_PREZ new.pptx"))
        pptx_txt = "\n".join(sh.text for sl in prs.slides for sh in sl.shapes if sh.has_text_frame)
        # graphiques : les valeurs sont dans les données de chart, pas le texte → on vérifie les cartes/textes
        contient(pptx_txt, ["0,918", "0,809", TESTS, auc_ext], "PowerPoint")
        contient(pptx_txt, ["693", "705", "720"], "PowerPoint", doit=False)
        if len(prs.slides._sldIdLst) < 20:
            err("PowerPoint : moins de 20 diapos (support incomplet)")
    except Exception as e:  # pptx non lisible = erreur de synchro
        err(f"PowerPoint illisible : {e}")

    # ---------- 5. Notebook ----------
    import nbformat
    nb = nbformat.read(ROOT / "notebooks/notebook_application_rondol.ipynb", as_version=4)
    nb_out = ""
    for c in nb.cells:
        if c.cell_type == "code":
            for o in c.get("outputs", []):
                if o.output_type == "error":
                    err(f"notebook : cellule en erreur ({o.ename})")
                nb_out += str(o.get("text", "")) + str(o.get("data", {}).get("text/plain", ""))
        else:
            nb_out += c.source
    if str(gen["n_rows"]) not in nb_out.replace(" ", "").replace(" ", ""):
        err("notebook : la sortie ne reflète pas la taille actuelle de la base consolidée")
    contient(nb_out, [f"{ext['roc_auc']}"], "notebook (validation externe)")

    # ---------- 6. Dossier de soutenance ----------
    # Le contrôle porte sur les sources Markdown, versionnées : un .docx binaire
    # ne se relit pas dans un diff, et c'est ce qui avait laissé « 705 tests »
    # survivre dans les guides précédents.
    src_dir = ROOT / "reports/soutenance/DOSSIER_FINAL/_source"
    dossier = " ".join(p.read_text(encoding="utf-8") for p in sorted(src_dir.glob("*.md")))
    if not dossier:
        err("dossier de soutenance : aucune source trouvée dans DOSSIER_FINAL/_source")
    # 0,809 est le F1-macro fold-aware retenu ; f1_rf (0,917) est celui du split
    # de test standard, il n'a pas sa place dans un support oral.
    contient(dossier, [acc_rf, auc_ext, TESTS, "100 800", "0,809",
                       "eu de rôle", "ntretien professionnel"], "dossier soutenance")
    contient(dossier, ["693", "705", "720 tests", "75 fichiers"],
             "dossier soutenance", doit=False)

    # ---------- 7. livrables binaires à jour (dates) ----------
    src_t = (ROOT / "docs/memoire_these_professionnelle_rondol.md").stat().st_mtime
    for f in ("MBEUMI_Wilfried_THESE.docx", "MBEUMI_Wilfried_THESE.pdf"):
        if (ROOT / f).stat().st_mtime < src_t:
            err(f"{f} plus ancien que le markdown source → rebuild requis")
    # Les PDF composés du dossier doivent suivre leurs sources Markdown.
    dossier_dir = ROOT / "reports/soutenance/DOSSIER_FINAL"
    for md in sorted((dossier_dir / "_source").glob("*.md")):
        pdf = next((p for p in dossier_dir.glob("*.pdf")
                    if p.stem.split(" - ")[0] == md.name[1]), None)
        if pdf and pdf.stat().st_mtime < md.stat().st_mtime - 5:
            err(f"{pdf.name} plus ancien que {md.name} → "
                f"relancer scripts/build_dossier_pdf.py")

    # ---------- verdict ----------
    if ERREURS:
        print("\n=== ÉCARTS DE SYNCHRONISATION ===")
        for e in ERREURS:
            print(" ✗", e)
        return 1
    print("\nSYNCHRONISATION OK — tous les livrables portent les mêmes chiffres.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
