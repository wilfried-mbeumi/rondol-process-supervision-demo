# -*- coding: utf-8 -*-
"""Soutenance PARTIE 1 (diapos 1–10) — design professionnel.

Usage : python scripts/build_soutenance_part1.py
Fusion : ouvrir les deux PPTX et copier les diapos de la partie 2 après la 10.
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "memoire_figures"
CAP = ROOT / "figures_memoire"
OUT = ROOT / "reports" / "soutenance" / "PREZ_part1.pptx"

# ── Palette ──────────────────────────────────────────────────────────────── #
BLEU       = RGBColor(0x1B, 0x3A, 0x5C)
BLEU_MED   = RGBColor(0x2E, 0x75, 0xB6)
BLEU_CLAIR = RGBColor(0xD6, 0xE4, 0xF0)
VERT       = RGBColor(0x1B, 0x7A, 0x3D)
VERT_CLAIR = RGBColor(0xD5, 0xF0, 0xDF)
ROUGE      = RGBColor(0xB0, 0x3A, 0x2E)
ROUGE_CLAIR= RGBColor(0xF9, 0xE2, 0xDF)
AMBRE      = RGBColor(0xB7, 0x79, 0x1F)
GRIS       = RGBColor(0x6B, 0x6B, 0x6B)
GRIS_CLAIR = RGBColor(0xF2, 0xF2, 0xF2)
ENCRE      = RGBColor(0x1E, 0x1E, 0x1E)
BLANC      = RGBColor(0xFF, 0xFF, 0xFF)
FOND       = RGBColor(0xFD, 0xFD, 0xFD)

POLICE = "Calibri"
L, H = Inches(13.333), Inches(7.5)
MARGE = Inches(0.75)
MARGE_D = Inches(0.75)
UTILE = L - MARGE - MARGE_D

# ── Helpers ──────────────────────────────────────────────────────────────── #

def _rect(diapo, x, y, w, h, fill=None, border=None, border_w=Pt(1), radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = diapo.shapes.add_shape(shape_type, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = border_w
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def _txt(diapo, x, y, w, h, texte, taille=16, gras=False, couleur=ENCRE,
         align=PP_ALIGN.LEFT, italique=False, interligne=1.3, ancre=None):
    zone = diapo.shapes.add_textbox(x, y, w, h)
    tf = zone.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = 0
    if ancre:
        tf.vertical_anchor = ancre
    p = tf.paragraphs[0]
    p.alignment = align; p.line_spacing = interligne
    r = p.add_run()
    r.text = texte; r.font.size = Pt(taille); r.font.bold = gras
    r.font.italic = italique; r.font.color.rgb = couleur; r.font.name = POLICE
    return zone, tf


def _add_para(tf, texte, taille=16, gras=False, couleur=ENCRE, avant=6,
              italique=False, interligne=1.3):
    p = tf.add_paragraph()
    p.space_before = Pt(avant); p.line_spacing = interligne
    r = p.add_run()
    r.text = texte; r.font.size = Pt(taille); r.font.bold = gras
    r.font.italic = italique; r.font.color.rgb = couleur; r.font.name = POLICE
    return p


def slide(prs):
    d = prs.slides.add_slide(prs.slide_layouts[6])
    d.background.fill.solid(); d.background.fill.fore_color.rgb = FOND
    return d


def header(d, titre, sous_titre=None, accent=BLEU):
    _rect(d, Inches(0), Inches(0), L, Inches(1.35), fill=accent)
    _rect(d, Inches(0), Inches(1.35), L, Inches(0.06), fill=VERT)
    _txt(d, MARGE, Inches(0.22), UTILE, Inches(0.85), titre,
         taille=28, gras=True, couleur=BLANC, interligne=1.05)
    if sous_titre:
        _txt(d, MARGE, Inches(1.5), UTILE, Inches(0.4), sous_titre,
             taille=13, couleur=GRIS, italique=True)


def notes(d, txt):
    d.notes_slide.notes_text_frame.text = txt.strip()


def kpi_card(d, x, y, w, h, valeur, label, couleur=BLEU, fond=BLANC, border=None):
    _rect(d, x, y, w, h, fill=fond, border=border or couleur, border_w=Pt(1.5), radius=True)
    _txt(d, x, y + Inches(0.18), w, Inches(0.7), valeur,
         taille=44, gras=True, couleur=couleur, align=PP_ALIGN.CENTER, interligne=1.0)
    _txt(d, x + Inches(0.1), y + Inches(0.95), w - Inches(0.2), Inches(0.9), label,
         taille=12, couleur=GRIS, align=PP_ALIGN.CENTER, interligne=1.25)


def kpi_row(d, items, y=Inches(2.3), couleur=BLEU, h=Inches(1.7)):
    n = len(items)
    gap = Inches(0.3)
    w = Emu(int((UTILE - gap * (n - 1)) / n))
    for i, (val, lab) in enumerate(items):
        x = Emu(int(MARGE + i * (w + gap)))
        kpi_card(d, x, y, w, h, val, lab, couleur=couleur)


def content_card(d, x, y, w, h, fill=BLANC, border=None):
    _rect(d, x, y, w, h, fill=fill, border=border, border_w=Pt(1), radius=True)


def blockquote(d, texte, y=Inches(5.5), couleur=BLEU, w=None):
    w = w or UTILE
    _rect(d, MARGE, y, Inches(0.07), Inches(1.2), fill=couleur)
    content_card(d, MARGE, y, w, Inches(1.2), fill=BLEU_CLAIR if couleur == BLEU else
                 VERT_CLAIR if couleur == VERT else ROUGE_CLAIR, border=couleur)
    _txt(d, MARGE + Inches(0.3), y + Inches(0.15), w - Inches(0.5), Inches(0.9),
         texte, taille=15, gras=True, couleur=couleur, interligne=1.3)


def bullet_block(d, items, y=Inches(2.2), taille=18, w=None, x=None):
    x = x or MARGE
    w = w or UTILE
    cur_y = y
    for item in items:
        if isinstance(item, tuple):
            titre, corps = item
        else:
            titre, corps = None, item
        card_h = Inches(0.85) if len(corps) < 120 else Inches(1.1)
        content_card(d, x, cur_y, w, card_h, fill=BLANC, border=GRIS_CLAIR)
        _rect(d, x, cur_y, Inches(0.06), card_h, fill=BLEU)
        text_x = x + Inches(0.22)
        text_w = w - Inches(0.35)
        if titre:
            _, tf = _txt(d, text_x, cur_y + Inches(0.1), text_w, Inches(0.35),
                         titre, taille=taille, gras=True, couleur=BLEU)
            _add_para(tf, corps, taille=taille - 2, couleur=ENCRE, avant=2)
        else:
            _txt(d, text_x, cur_y + Inches(0.12), text_w, card_h - Inches(0.2),
                 corps, taille=taille - 1, couleur=ENCRE)
        cur_y += card_h + Inches(0.12)


def fig(d, chemin, haut=Inches(2.0), bas=Inches(0.3)):
    from PIL import Image as PILImage
    iw, ih = PILImage.open(chemin).size
    dispo_h = H - haut - bas
    dispo_w = UTILE
    ratio = iw / ih
    w = dispo_w; h = Emu(int(w / ratio))
    if h > dispo_h:
        h = dispo_h; w = Emu(int(h * ratio))
    content_card(d, Emu(int((L - w) / 2)) - Inches(0.1), haut - Inches(0.08),
                 Emu(int(w)) + Inches(0.2), Emu(int(h)) + Inches(0.16),
                 fill=BLANC, border=GRIS_CLAIR)
    d.shapes.add_picture(str(chemin), Emu(int((L - w) / 2)), haut,
                         width=Emu(int(w)), height=Emu(int(h)))


def num_slide(d, numero, total=20):
    _txt(d, L - Inches(1.0), H - Inches(0.4), Inches(0.7), Inches(0.3),
         f"{numero}/{total}", taille=10, couleur=GRIS, align=PP_ALIGN.RIGHT)


# =========================================================================== #
prs = Presentation()
prs.slide_width, prs.slide_height = L, H

# ── 1. TITRE ──────────────────────────────────────────────────────────────── #
d = slide(prs)
d.background.fill.solid(); d.background.fill.fore_color.rgb = BLEU
_rect(d, Inches(0), Inches(0), Inches(0.35), H, fill=VERT)
_rect(d, Inches(0), H - Inches(0.08), L, Inches(0.08), fill=VERT)

_txt(d, Inches(1.2), Inches(1.2), Inches(10.8), Inches(0.35),
     "THÈSE PROFESSIONNELLE · MASTÈRE 2 DATA & INTELLIGENCE ARTIFICIELLE",
     taille=12, gras=True, couleur=BLEU_CLAIR, interligne=1.0)

_rect(d, Inches(1.2), Inches(1.7), Inches(2.5), Inches(0.04), fill=VERT)

_txt(d, Inches(1.2), Inches(2.0), Inches(10.8), Inches(1.8),
     "Rendre un procédé d'extrusion\nlisible, comparable et prédictible",
     taille=42, gras=True, couleur=BLANC, interligne=1.15)

_txt(d, Inches(1.2), Inches(4.1), Inches(10.8), Inches(0.7),
     "Jumeau numérique et IA explicable pour l'extrusion bivis\n"
     "de composants de batteries tout-solide",
     taille=17, couleur=BLEU_CLAIR, interligne=1.35)

_rect(d, Inches(1.2), Inches(5.3), Inches(8.0), Inches(0.01), fill=RGBColor(0x4A, 0x7A, 0xA5))

_txt(d, Inches(1.2), Inches(5.6), Inches(10.8), Inches(1.0),
     "Wilfried Galtier MBEUMI\n"
     "Nexa Digital School   ·   Encadrement industriel : Maël Gallas — Rondol Industrie   ·   2025–2026",
     taille=14, couleur=BLANC, interligne=1.55)
notes(d, """Bonjour. Je m'appelle Wilfried MBEUMI et je vais vous présenter le travail
mené chez Rondol Industrie pendant quatre mois.
Le titre annonce l'ambition : rendre un procédé d'extrusion lisible, comparable et
prédictible. Ces trois mots structureront toute ma présentation.
[~30 secondes]""")

# ── 2. LE VERROU ──────────────────────────────────────────────────────────── #
d = slide(prs)
header(d, "Le verrou : on sait faire la chimie, pas la mise en forme",
       "Batteries tout-solide — un blocage devenu procédé")
num_slide(d, 2)

bullet_block(d, [
    ("Le blocage n'est plus chimique.",
     "Les batteries tout-solide promettent plus de densité d'énergie et moins de risque "
     "d'incendie. Ce qui bloque aujourd'hui, c'est la mise en forme des électrodes."),
    ("La voie humide domine encore.",
     "Elle impose la NMP, un solvant classé toxique pour la reproduction, et un séchage énergivore."),
    ("L'extrusion bivis sèche est l'alternative.",
     "Continue, sans solvant — et c'est le savoir-faire historique de Rondol. "
     "Mais elle reste très peu documentée pour ces formulations céramiques abrasives."),
], y=Inches(2.1), taille=18)

content_card(d, MARGE, Inches(5.9), UTILE, Inches(0.7), fill=GRIS_CLAIR)
_txt(d, MARGE + Inches(0.15), Inches(5.98), UTILE - Inches(0.3), Inches(0.55),
     "Marché aval : 0,26 → 1,77 Md USD d'ici 2031 (TCAC 37,5 %, MarketsandMarkets) ; "
     "Grand View Research retient 1,60 → 15,65 Md USD (TCAC 31,8 %).",
     taille=12, couleur=GRIS, italique=True)
notes(d, """Le contexte en trois temps.
D'abord : la difficulté n'est plus de formuler la chimie, mais de mettre en forme
l'électrode. C'est un problème de procédé.
Ensuite : la voie dominante utilise la NMP, un solvant toxique pour la reproduction.
Enfin : l'extrusion sèche est l'alternative, et c'est justement le métier de Rondol.
[~2 minutes]""")

# ── 3. LE PROBLÈME CONCRET ─────────────────────────────────────────────── #
d = slide(prs)
header(d, "Chez Rondol, le problème est concret",
       "Extrudeuse bivis Ø 10,5 mm — R&D lithium")
num_slide(d, 3)

kpi_row(d, [("81", "positions de vis\nà configurer"),
            ("13", "types d'éléments\ndisponibles"),
            ("12", "capteurs de\ntempérature"),
            ("8", "essais\nexploitables")],
        y=Inches(2.3))

blockquote(d, "L'opérateur règle par expérience. Rien ne compare deux configurations. "
              "Rien n'anticipe une dérive.", y=Inches(4.6), couleur=ROUGE)

content_card(d, MARGE, Inches(6.2), UTILE, Inches(0.7), fill=GRIS_CLAIR)
_txt(d, MARGE + Inches(0.15), Inches(6.28), UTILE - Inches(0.3), Inches(0.55),
     "Campagne d'essais du 7 au 13 avril 2026 — la seule base expérimentale disponible.",
     taille=12, couleur=GRIS, italique=True, align=PP_ALIGN.CENTER)
notes(d, """Voici le problème tel que je l'ai trouvé en arrivant.
81 positions de vis : aucun opérateur ne peut explorer cet espace à la main.
12 capteurs qui produisent des données en continu — mais ces données servent à
surveiller après coup, jamais à anticiper.
Et 8 essais exploitables seulement. Retenez ce chiffre : il va conditionner toute
la partie modélisation.
[~2 minutes]""")

# ── 4. LA PROBLÉMATIQUE ──────────────────────────────────────────────────── #
d = slide(prs)
header(d, "La problématique")
num_slide(d, 4)

content_card(d, MARGE, Inches(2.0), UTILE, Inches(1.6), fill=BLEU_CLAIR, border=BLEU)
_rect(d, MARGE, Inches(2.0), Inches(0.07), Inches(1.6), fill=BLEU)
_txt(d, MARGE + Inches(0.35), Inches(2.15), UTILE - Inches(0.6), Inches(1.3),
     "Comment concevoir un système d'aide à la décision qui rende un procédé "
     "d'extrusion bivis lisible, comparable et prédictible — sans jamais "
     "devenir une boîte noire ?",
     taille=22, gras=True, couleur=BLEU, interligne=1.3)

_txt(d, MARGE, Inches(3.9), UTILE, Inches(0.35),
     "Trois exigences non négociables, fixées dès le cadrage :",
     taille=14, couleur=GRIS, italique=True)

exigences = [("Traçable", "Toute valeur affichée doit pouvoir être\nexpliquée à un ingénieur procédé", BLEU),
             ("Honnête", "Ce qui n'est pas calibré industriellement\nest annoncé comme tel", VERT),
             ("Démontrable", "Un outil réellement utilisable devant\nle client, pas une maquette", AMBRE)]
gap = Inches(0.3)
w = Emu(int((UTILE - gap * 2) / 3))
for i, (titre, desc, coul) in enumerate(exigences):
    x = Emu(int(MARGE + i * (w + gap)))
    content_card(d, x, Inches(4.45), w, Inches(2.1), fill=BLANC, border=coul)
    _rect(d, x, Inches(4.45), w, Inches(0.5), fill=coul)
    _txt(d, x, Inches(4.5), w, Inches(0.4), titre,
         taille=18, gras=True, couleur=BLANC, align=PP_ALIGN.CENTER)
    _txt(d, x + Inches(0.15), Inches(5.1), w - Inches(0.3), Inches(1.3), desc,
         taille=14, couleur=ENCRE, align=PP_ALIGN.CENTER, interligne=1.3)

notes(d, """Voici la question à laquelle ce travail répond.
Insistez sur la fin : « sans jamais devenir une boîte noire ».
Les trois exigences ont été posées au cadrage avec Maël Gallas et n'ont pas bougé.
La deuxième — l'honnêteté — reviendra de façon décisive en seconde partie.
[~1 minute 30]""")

# ── 5. ARCHITECTURE ──────────────────────────────────────────────────────── #
d = slide(prs)
header(d, "Architecture en couches : envelopper, ne pas recalculer",
       "Le calcul procédé est appelé une seule fois ; tout le reste réutilise son résultat")
num_slide(d, 5)

fig(d, FIG / "fig_architecture.png", haut=Inches(2.1))

notes(d, """L'architecture en cinq couches, avec un principe fondateur : envelopper,
ne pas recalculer.
Le calcul procédé — la géométrie de vis — est appelé UNE SEULE FOIS. Toutes les
couches supérieures réutilisent son résultat au lieu de le recalculer.
Pourquoi c'est important : c'est ce qui garantit qu'aucun chiffre affiché ne peut
contredire un chiffre affiché ailleurs.
[~2 minutes]""")

# ── 6. LE CŒUR MÉTIER ────────────────────────────────────────────────────── #
d = slide(prs)
header(d, "Le cœur métier : la vis, modélisée position par position",
       "Convoyage → malaxage / cisaillement → restriction → pompage filière")
num_slide(d, 6)

fig(d, FIG / "fig_engine_logic.png", haut=Inches(2.1), bas=Inches(1.6))

content_card(d, MARGE, Inches(6.05), UTILE, Inches(0.85), fill=GRIS_CLAIR)
_txt(d, MARGE + Inches(0.15), Inches(6.12), UTILE - Inches(0.3), Inches(0.7),
     "Ce n'est pas un modèle générique d'extrusion : c'est la géométrie réelle de "
     "l'extrudeuse Ø 10,5 mm de Rondol, élément par élément. Chaque configuration "
     "devient un objet calculable, comparable et historisable.",
     taille=13, couleur=GRIS, italique=True, interligne=1.3)
notes(d, """Le cœur du système, c'est la modélisation de la vis position par position.
De la géométrie réelle, le moteur dérive le taux de remplissage, le temps de séjour,
le volume occupé et le volume libre — puis les agrège par zone thermique.
Point clé : ce n'est pas un modèle générique. C'est la machine de Rondol.
[~2 minutes]""")

# ── 7. LA PHYSIQUE ────────────────────────────────────────────────────────── #
d = slide(prs)
header(d, "La physique embarquée — nominale, et annoncée comme telle",
       "Chaque équation est documentée ; ce qui n'est pas calibré est signalé")
num_slide(d, 7)

eqs = [("Viscosité locale", "Carreau-Yasuda couplé à Arrhenius,\npresets par matière (LFP, LATP, liants fluorés)", BLEU),
       ("Couple local", "M = η·γ̇²·V_rempli / (2πN)\ncalculé nœud par nœud puis agrégé", BLEU_MED),
       ("Équation thermique", "T_réel = T_consigne + (2πN·M)/(ṁ·Cp) + k·τ\nimposée par l'encadrement industriel", VERT)]
for i, (titre, desc, coul) in enumerate(eqs):
    y = Inches(2.1) + i * Inches(1.25)
    content_card(d, MARGE, y, UTILE, Inches(1.1), fill=BLANC, border=GRIS_CLAIR)
    _rect(d, MARGE, y, Inches(0.06), Inches(1.1), fill=coul)
    _txt(d, MARGE + Inches(0.25), y + Inches(0.12), Inches(3.5), Inches(0.35),
         titre, taille=17, gras=True, couleur=coul)
    _txt(d, MARGE + Inches(3.9), y + Inches(0.12), UTILE - Inches(4.1), Inches(0.85),
         desc, taille=14, couleur=ENCRE, interligne=1.3)

content_card(d, MARGE, Inches(5.9), UTILE, Inches(0.95), fill=ROUGE_CLAIR, border=ROUGE)
_txt(d, MARGE + Inches(0.2), Inches(5.98), UTILE - Inches(0.4), Inches(0.8),
     "Ce qui n'est pas fait est écrit : énergie mécanique locale, température avancée "
     "et pression filière restent des briques différées, affichées « À venir » "
     "dans l'interface plutôt que remplies de valeurs plausibles.",
     taille=13, couleur=ROUGE, italique=True, interligne=1.3)
notes(d, """Trois équations embarquées : viscosité, couple local, équation thermique.
Le point en bas : ce qui n'est pas calculé est affiché « À venir » dans l'interface.
J'aurais pu remplir ces cases avec des valeurs plausibles — j'ai préféré le vide honnête.
[~2 minutes]""")

# ── 8. L'AGENT ────────────────────────────────────────────────────────────── #
d = slide(prs)
header(d, "Un agent explicable, pas un classifieur déguisé",
       "État procédé → règles métier explicites → alertes → recommandations chiffrées")
num_slide(d, 8)

fig(d, FIG / "fig_two_level_ai.png", haut=Inches(2.1), bas=Inches(1.7))

content_card(d, MARGE, Inches(6.1), UTILE, Inches(0.85), fill=VERT_CLAIR, border=VERT)
_txt(d, MARGE + Inches(0.2), Inches(6.18), UTILE - Inches(0.4), Inches(0.7),
     "La décision n'est jamais confiée au modèle statistique : le modèle prédit "
     "un score, les règles recommandent, l'humain tranche.",
     taille=16, gras=True, couleur=VERT, align=PP_ALIGN.CENTER, interligne=1.3)
notes(d, """L'agent fonctionne par règles métier explicites, pas par apprentissage.
Chaque alerte cite sa preuve chiffrée. L'opérateur peut vérifier.
La séparation des rôles est nette : le modèle prédit, les règles recommandent, l'humain décide.
[~2 minutes]""")

# ── 9. LES DONNÉES ───────────────────────────────────────────────────────── #
d = slide(prs)
header(d, "Les données : une campagne réelle, et son revers",
       "Campagne Rondol du 7 au 13 avril 2026")
num_slide(d, 9)

kpi_row(d, [("310 782", "relevés bruts\n12 capteurs · 25,8 Mo"),
            ("10–16 %", "couverture par capteur\nacquisition fragmentée"),
            ("627", "fenêtres de 60 s\nexploitables"),
            ("87", "variables extraites\n12 × 7 stats + 3 croisées")],
        y=Inches(2.2), h=Inches(1.65))

blockquote(d, "Huit essais. C'est peu. Je l'ai traité comme la limite du projet, "
              "pas comme un détail à contourner.", y=Inches(4.5), couleur=AMBRE)

content_card(d, MARGE, Inches(6.1), UTILE, Inches(0.7), fill=GRIS_CLAIR)
_txt(d, MARGE + Inches(0.15), Inches(6.18), UTILE - Inches(0.3), Inches(0.55),
     "Codes d'erreur thermocouple à 3276,7 °C · doublons temporels · "
     "couverture capteur hétérogène — nettoyés et documentés.",
     taille=12, couleur=GRIS, italique=True, align=PP_ALIGN.CENTER)
notes(d, """Les données viennent d'une campagne réelle, pas d'une simulation.
Mais regardez la deuxième colonne : 10 à 16 % de couverture seulement.
Après nettoyage il reste 627 fenêtres exploitables, issues de 8 essais.
Cette phrase en bas est importante : je n'ai pas cherché à masquer la faiblesse du volume.
[~2 minutes]""")

# ── 10. LE PIPELINE ML ───────────────────────────────────────────────────── #
d = slide(prs)
header(d, "Le pipeline ML — la rigueur avant la performance",
       "Séparation par essai · cible décalée · Leave-One-Group-Out")
num_slide(d, 10)

fig(d, FIG / "fig_validation.png", haut=Inches(2.1), bas=Inches(2.2))

content_card(d, MARGE, Inches(5.3), Inches(5.5), Inches(1.5), fill=BLANC, border=BLEU)
_txt(d, MARGE + Inches(0.2), Inches(5.35), Inches(5.1), Inches(0.25),
     "87,5 % d'entraînement par pli en moyenne", taille=13, gras=True, couleur=BLEU)
_txt(d, MARGE + Inches(0.2), Inches(5.65), Inches(5.1), Inches(1.0),
     "Partition naïve :  F1 = 0,92\n"
     "Séparation stricte :  F1 = 0,79\n"
     "→ Quinze points d'écart", taille=15, couleur=ENCRE, interligne=1.5)

content_card(d, MARGE + Inches(6.0), Inches(5.3), UTILE - Inches(6.0), Inches(1.5),
             fill=BLEU_CLAIR, border=BLEU)
_txt(d, MARGE + Inches(6.2), Inches(5.4), UTILE - Inches(6.4), Inches(1.3),
     "Ce n'est pas un défaut à cacher : c'est la mesure de ce que vaut vraiment le modèle.",
     taille=15, gras=True, couleur=BLEU, interligne=1.3)

notes(d, """Le protocole de validation est le point méthodologique le plus important.
Les fenêtres d'un même essai sont fortement autocorrélées. Les répartir au hasard
ferait fuir l'information et gonflerait les scores.
D'où la séparation par essai et le Leave-One-Group-Out : chaque essai est écarté à
son tour, l'entraînement porte sur 87,5 % des fenêtres en moyenne.
Le résultat : 0,92 en naïve, 0,79 en stricte. Quinze points.
[~2 minutes 30]""")

# ── Sauvegarde ────────────────────────────────────────────────────────────── #
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"[OK] {OUT}  —  {sum(1 for _ in prs.slides)} diapositive(s)")
